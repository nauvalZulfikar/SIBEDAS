"""
Build profesional proposal docs (.docx, .pdf, .pptx) untuk fitur Monitoring Satelit.
"""
import os
from datetime import datetime

# ====================== CONSTANTS ======================
OUT_DIR = os.path.dirname(os.path.abspath(__file__))
CLIENT_NAME = "DPUTR Kab. Bandung"
PROJECT_NAME = "Sistem Monitoring Satelit Bangunan"
VENDOR_NAME = "Nauval Zulfikar"
DATE_STR = datetime.now().strftime("%d %B %Y")
PROPOSAL_NO = f"PROP-SATELIT-{datetime.now().strftime('%Y%m')}-001"

PRICE_TOTAL = 150_000_000
PRICE_A_TRANSFER = 15_000_000
PRICE_B_SATELIT = 120_000_000
PRICE_C_TRAINING = 8_000_000
PRICE_D_MAINT = 7_000_000

# Colors (brand)
NAVY = (28, 45, 84)       # primary
GOLD = (201, 162, 39)     # accent
GRAY = (90, 90, 90)
LIGHT = (245, 247, 250)
DANGER = (220, 53, 69)
SUCCESS = (40, 167, 69)

# ====================== DOCX ======================
def build_docx():
    from docx import Document
    from docx.shared import Pt, RGBColor, Cm, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
    from docx.enum.table import WD_ALIGN_VERTICAL
    from docx.oxml.ns import qn, nsmap
    from docx.oxml import OxmlElement

    doc = Document()
    # Page setup
    for s in doc.sections:
        s.top_margin = Cm(2.2); s.bottom_margin = Cm(2.2); s.left_margin = Cm(2.5); s.right_margin = Cm(2.0)

    styles = doc.styles
    # Normal style
    n = styles['Normal']
    n.font.name = 'Calibri'; n.font.size = Pt(11)

    def set_cell_bg(cell, rgb_hex):
        tc = cell._tc; tcPr = tc.get_or_add_tcPr()
        shd = OxmlElement('w:shd')
        shd.set(qn('w:val'), 'clear'); shd.set(qn('w:color'), 'auto'); shd.set(qn('w:fill'), rgb_hex)
        tcPr.append(shd)

    def H(text, level=1, color=NAVY, size=None, align=None, bold=True):
        sizes = {1: 22, 2: 16, 3: 13}
        p = doc.add_paragraph()
        if align == 'center': p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r = p.add_run(text)
        r.font.size = Pt(size or sizes.get(level, 13))
        r.font.bold = bold
        r.font.color.rgb = RGBColor(*color)
        return p

    def P(text, bold=False, italic=False, color=None, size=11, align=None, space_after=6):
        p = doc.add_paragraph()
        if align == 'center': p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_after = Pt(space_after)
        r = p.add_run(text)
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        if color: r.font.color.rgb = RGBColor(*color)
        return p

    def rupiah(n):
        return f"Rp {n:,.0f}".replace(",", ".")

    # ============ COVER ============
    for _ in range(3): doc.add_paragraph()
    H("PROPOSAL PENAWARAN", level=1, size=36, align='center', color=NAVY)
    H(PROJECT_NAME, level=2, size=20, align='center', color=GOLD)
    doc.add_paragraph()
    P(f"Diajukan kepada:", align='center', color=GRAY, size=11)
    P(CLIENT_NAME, align='center', bold=True, size=14, color=NAVY)
    doc.add_paragraph(); doc.add_paragraph()

    tbl = doc.add_table(rows=3, cols=2); tbl.autofit = False
    tbl.columns[0].width = Cm(5); tbl.columns[1].width = Cm(10)
    for (lbl, val) in [("Nomor Proposal", PROPOSAL_NO), ("Tanggal", DATE_STR), ("Penyedia", VENDOR_NAME)]:
        row = tbl.add_row() if tbl.rows else None
    for i, (lbl, val) in enumerate([("Nomor Proposal", PROPOSAL_NO), ("Tanggal", DATE_STR), ("Penyedia", VENDOR_NAME)]):
        tbl.cell(i, 0).text = lbl; tbl.cell(i, 1).text = val
        set_cell_bg(tbl.cell(i, 0), 'E8ECEF')

    doc.add_paragraph(); doc.add_paragraph()
    P("“Fitur ini membayar dirinya sendiri kalau DPUTR berhasil", align='center', italic=True, color=NAVY, size=12)
    P("menertibkan 15 bangunan saja dari 96.322 yang terdeteksi.”", align='center', italic=True, color=NAVY, size=12)
    doc.add_page_break()

    # ============ EXECUTIVE SUMMARY ============
    H("RINGKASAN EKSEKUTIF", level=1)
    P("Kabupaten Bandung memiliki ribuan bangunan yang belum memiliki Persetujuan Bangunan Gedung (PBG) — sumber potensi pendapatan asli daerah (PAD) yang selama ini tidak termonitor. Sistem Monitoring Satelit ini memanfaatkan citra satelit + data SIMBG untuk mendeteksi, mengklasifikasi, dan memetakan seluruh bangunan di 31 kecamatan Kabupaten Bandung secara otomatis.")
    doc.add_paragraph()
    H("Angka Kunci", level=2)
    t = doc.add_table(rows=3, cols=2); t.style = 'Light Grid Accent 1'
    kv = [
        ("Total bangunan terdeteksi", "1.09 juta"),
        ("Bangunan ≥200 m² tanpa izin sah", "96.322"),
        ("Potensi PAD (skenario 10%)", "Rp 102,7 miliar"),
    ]
    for i, (k, v) in enumerate(kv):
        t.cell(i,0).text = k; t.cell(i,1).text = v
        for c in t.rows[i].cells:
            c.vertical_alignment = WD_ALIGN_VERTICAL.CENTER
        set_cell_bg(t.cell(i,0), 'E8ECEF')
    doc.add_paragraph()
    P("Dari 5.241 PBG yang sudah terbit di Kab. Bandung, rata-rata retribusi yang masuk = Rp 10,66 juta per PBG.", italic=True, color=GRAY)
    P("Artinya: 96.322 bangunan yang belum tertangani setara dengan Rp 1,03 triliun potensi PAD yang belum dieksekusi.", bold=True, color=SUCCESS, size=12)
    doc.add_page_break()

    # ============ 1. LATAR BELAKANG ============
    H("1. LATAR BELAKANG", level=1)
    P("DPUTR Kabupaten Bandung membutuhkan instrumen monitoring yang dapat:")
    for b in [
        "Mendeteksi bangunan tanpa izin di 31 kecamatan secara otomatis.",
        "Menyajikan data visual peta interaktif untuk staff enforcement.",
        "Mencocokkan data deteksi dengan record PBG SIMBG existing.",
        "Memberikan dasar kuantitatif untuk target PAD tahunan.",
    ]:
        doc.add_paragraph(b, style='List Bullet')
    doc.add_paragraph()
    P("Metode konvensional (survey lapangan) memakan waktu berbulan-bulan dan biaya tinggi. Solusi berbasis satelit + SIMBG menghasilkan data yang sama dalam hitungan menit dan bisa di-refresh kapan saja.")
    doc.add_page_break()

    # ============ 2. OUTPUT & OUTCOME ============
    H("2. OUTPUT & OUTCOME", level=1)
    H("Output — Yang DPUTR Terima", level=2)
    for b in [
        "Aplikasi web Monitoring Satelit yang terdeploy di production server DPUTR.",
        "Peta interaktif 31 polygon kecamatan dengan 97.000+ titik deteksi bangunan.",
        "Dashboard KPI realtime dan tabel breakdown per kecamatan.",
        "Workflow verifikasi untuk staff enforcement dengan audit trail.",
        "Source code + database schema + dokumentasi teknis & user manual.",
        "1 batch training untuk staff DPUTR.",
    ]:
        doc.add_paragraph(b, style='List Bullet')

    H("Outcome — Manfaat Terukur", level=2)
    for b in [
        "Visibilitas 96.322 bangunan potensial wajib PBG yang selama ini tidak termonitor.",
        "Efisiensi enforcement: target bangunan tanpa izin langsung dipetakan, tidak perlu survey acak.",
        "Basis data terpadu untuk sinkronisasi dengan PBB, tata ruang, dan perizinan lainnya.",
        "Data-driven decision: Bupati & DPRD punya angka konkret untuk kebijakan PAD.",
        "Dokumentasi kontribusi kinerja DPUTR dalam laporan tahunan.",
    ]:
        doc.add_paragraph(b, style='List Bullet')
    doc.add_page_break()

    # ============ 3. POTENSI PENDAPATAN ============
    H("3. POTENSI PENDAPATAN (ROI)", level=1)
    P("Berdasarkan data real SIMBG Kab. Bandung (5.241 PBG yang sudah terbit):", italic=True)
    t = doc.add_table(rows=2, cols=2); t.style = 'Light Grid Accent 1'
    t.cell(0,0).text = "Rata-rata retribusi per PBG"; t.cell(0,1).text = "Rp 10.660.000"
    t.cell(1,0).text = "Total PAD PBG terkumpul"; t.cell(1,1).text = "Rp 55,86 miliar"
    for r in t.rows: set_cell_bg(r.cells[0], 'E8ECEF')
    doc.add_paragraph()

    H("Skenario Potensi Pendapatan dari 96.322 Bangunan Tanpa Izin", level=2)
    t = doc.add_table(rows=1, cols=3); t.style = 'Light Grid Accent 1'
    hdr = t.rows[0].cells
    hdr[0].text = "Skenario Enforcement"; hdr[1].text = "PBG Tertarik"; hdr[2].text = "Potensi PAD"
    for c in hdr: set_cell_bg(c, '1C2D54')
    for c in hdr:
        for p in c.paragraphs:
            for r in p.runs:
                r.font.color.rgb = RGBColor(255,255,255); r.font.bold = True
    rows_data = [
        ("1% (super konservatif)", "963", "Rp 10,3 miliar"),
        ("5% (konservatif)", "4.816", "Rp 51,3 miliar"),
        ("10% (realistis)", "9.632", "Rp 102,7 miliar"),
        ("20% (target wajar)", "19.264", "Rp 205,4 miliar"),
        ("30% (agresif)", "28.897", "Rp 308,0 miliar"),
    ]
    for rd in rows_data:
        row = t.add_row().cells
        for i, v in enumerate(rd): row[i].text = v

    doc.add_paragraph()
    P("Skenario realistis saja (10% enforcement) sudah setara membangun 3 jembatan baru.", italic=True, color=NAVY, size=12, align='center')
    doc.add_paragraph()

    H("Manfaat Tambahan (Belum Dihitung)", level=2)
    for b in [
        "Denda administratif bangunan tanpa izin: 2-10% nilai bangunan.",
        "Akurasi penagihan PBB: potensi tambahan Rp 10-20 miliar/tahun.",
        "Compliance tata ruang & RDTR enforcement.",
        "Decision support untuk perencanaan infrastruktur.",
    ]:
        doc.add_paragraph(b, style='List Bullet')
    doc.add_page_break()

    # ============ 4. RUANG LINGKUP ============
    H("4. RUANG LINGKUP PEKERJAAN", level=1)
    P("Seluruh scope pekerjaan disajikan dalam bentuk tabel untuk kemudahan audit dan tracking deliverable. Rincian nilai investasi ada pada Bagian 5.", italic=True, color=GRAY)
    doc.add_paragraph()

    def pkg_table_nopr(code, name, items):
        """Render satu paket sebagai tabel TANPA harga (harga dipisah ke section Investasi)."""
        rows = [["Kode Paket", code],
                ["Nama Paket", name]]
        t = doc.add_table(rows=len(rows) + 1 + len(items), cols=2)
        t.style = 'Light Grid Accent 1'
        t.columns[0].width = Cm(5); t.columns[1].width = Cm(11)
        for i, (k, v) in enumerate(rows):
            t.cell(i, 0).text = k; t.cell(i, 1).text = v
            set_cell_bg(t.cell(i, 0), '1C2D54')
            for p in t.cell(i, 0).paragraphs:
                for r in p.runs: r.font.color.rgb = RGBColor(255,255,255); r.font.bold = True
        hdr_idx = len(rows)
        merged = t.cell(hdr_idx, 0).merge(t.cell(hdr_idx, 1))
        merged.text = "Deliverable & Fitur Tercakup"
        set_cell_bg(merged, 'C9A227')
        for p in merged.paragraphs:
            for r in p.runs: r.font.color.rgb = RGBColor(255,255,255); r.font.bold = True
        for j, item in enumerate(items):
            row_idx = hdr_idx + 1 + j
            t.cell(row_idx, 0).text = f"{code}.{j+1}"
            t.cell(row_idx, 1).text = item
            set_cell_bg(t.cell(row_idx, 0), 'E8ECEF')
            for p in t.cell(row_idx, 0).paragraphs:
                for r in p.runs: r.font.bold = True; r.font.color.rgb = RGBColor(*NAVY)
        doc.add_paragraph()

    # ===== PAKET A =====
    H("Paket A — Transfer, Deployment & Rekonsiliasi", level=2)
    pkg_table_nopr("A", "Transfer, Deployment & Rekonsiliasi Data", [
        "Setup production environment (nginx, PHP-FPM, MariaDB tuning).",
        "Import & migrate database produksi (297 MB).",
        "Bug-fix login session, Vite pipeline, asset loading.",
        "Rekonsiliasi angka dengan Dashboard Pimpinan SIMBG & Dashboard PBG.",
        "Deployment runbook + rollback procedure.",
        "Handover documentation + 1 batch training.",
    ])

    # ===== PAKET B =====
    H("Paket B — Pembangunan Fitur Monitoring Satelit", level=2)
    P("Paket B dipecah menjadi 7 sub-paket teknis:", italic=True)

    subpkg = [
        ("B.1", "Data Pipeline & Integrasi", [
            "Ingest dataset Microsoft Building Footprints (1,09 juta bangunan).",
            "Integrasi data PBG SIMBG (pbg_task, pbg_task_details, retributions).",
            "Spatial matching engine deteksi satelit ↔ PBG.",
            "Integrasi polygon resmi 31 kecamatan Kab. Bandung (sumber BPS).",
            "Point-in-polygon classification 1,09M rows via Python shapely.",
            "Data quality audit (orphan FK, validasi spasial).",
        ]),
        ("B.2", "Backend API", [
            "7 REST endpoint: stats, geojson, pbg-geojson, verify, refresh.",
            "Multi-filter query (kecamatan, jenis, status, luas).",
            "2-tier caching (Laravel cache + DB snapshot).",
            "Indexing & query optimization (cold 15s → warm 0,4s).",
            "Authentication layer (Sanctum + role-based).",
        ]),
        ("B.3", "Frontend Peta Interaktif", [
            "Map engine Leaflet + Esri satellite tile + Carto labels.",
            "Bounded scroll + dynamic min-zoom (area lock Kab. Bandung).",
            "Render 31 polygon kecamatan + click-to-zoom + tooltip.",
            "Marker clustering 2-layer (deteksi satelit + titik PBG).",
            "Popup detail bangunan + panel verifikasi.",
            "Loading indicator + abort-on-pan race handling.",
            "Legend 5-state (SK Terbit / Proses / Ditolak / Tanpa Izin / PBG).",
            "Filter system 6-kategori (kecamatan, jenis, status, luas, tampilan, layer).",
        ]),
        ("B.4", "Dashboard Analytics", [
            "4 KPI cards realtime (Total / Berizin / Tanpa Izin / Rasio).",
            "PBG breakdown strip (Total / Terbit / Proses / Ditolak).",
            "Tabel 'Data per Kecamatan' dengan stacked bar 3-segment.",
            "Grand total row + snapshot timestamp.",
            "Breakdown jenis bangunan (7 kategori fungsi).",
            "Drill-down interaktif (klik row → zoom + filter).",
        ]),
        ("B.5", "Verification & Audit Workflow", [
            "Panel verifikasi 4-status (Illegal / Legal / Review / False+).",
            "Bulk verification API + audit trail (verified_by, verified_at, notes).",
            "Role-based access control.",
        ]),
        ("B.6", "Performance & Production Readiness", [
            "Tabel snapshot kecamatan_stats (pre-computed, 186 rows).",
            "Artisan command kecamatan-stats:refresh (idempotent, ~60s).",
            "Cache strategy dengan TTL + auto-invalidation on verify.",
        ]),
        ("B.7", "Dokumentasi & Proof-of-Correctness", [
            "SQL audit script verifikasi angka.",
            "User manual staff PUTR (screenshot + workflow).",
            "Technical documentation (schema ERD, API spec).",
        ]),
    ]
    total_rows = 1 + sum(1 + len(items) for _, _, items in subpkg)
    tb = doc.add_table(rows=total_rows, cols=3); tb.style = 'Light Grid Accent 1'
    tb.columns[0].width = Cm(2); tb.columns[1].width = Cm(5); tb.columns[2].width = Cm(9)
    hdr = tb.rows[0].cells
    hdr[0].text = "Kode"; hdr[1].text = "Sub-Paket"; hdr[2].text = "Deliverable"
    for c in hdr:
        set_cell_bg(c, '1C2D54')
        for p in c.paragraphs:
            for r in p.runs: r.font.color.rgb = RGBColor(255,255,255); r.font.bold = True
    ri = 1
    for code, name, items in subpkg:
        row = tb.rows[ri].cells
        row[0].text = code; row[1].text = name; row[2].text = ""
        for c in row:
            set_cell_bg(c, 'C9A227')
            for p in c.paragraphs:
                for r in p.runs: r.font.bold = True; r.font.color.rgb = RGBColor(255,255,255)
        ri += 1
        for j, item in enumerate(items):
            row = tb.rows[ri].cells
            row[0].text = f"{code}.{j+1}"; row[1].text = ""; row[2].text = item
            set_cell_bg(row[0], 'E8ECEF')
            for p in row[0].paragraphs:
                for r in p.runs: r.font.bold = True; r.font.color.rgb = RGBColor(*NAVY)
            ri += 1
    doc.add_paragraph()

    # ===== PAKET C =====
    H("Paket C — Data Enrichment Geocoding (Opsional)", level=2)
    pkg_table_nopr("C", "Enrichment 96.322 Bangunan dengan Alamat Lengkap", [
        "Reverse geocoding seluruh koordinat deteksi ke alamat (jalan, RT/RW, kelurahan).",
        "Integrasi Badan Informasi Geospasial (BIG) Tanah Air sebagai source utama (gratis).",
        "Fallback Nominatim OpenStreetMap untuk area yang tidak tercover BIG.",
        "Batch processing 96k+ row dengan queue + retry handling.",
        "Cache hasil geocoding di DB untuk efisiensi query lanjutan.",
        "Tampil di popup marker + export CSV lengkap alamat.",
        "Export ready-to-use untuk surat teguran (template mail-merge).",
    ])

    # ===== PAKET D =====
    H("Paket D — Maintenance & Support 6 Bulan (Opsional)", level=2)
    pkg_table_nopr("D", "Maintenance & Support 6 Bulan", [
        "Bug-fix unlimited untuk defect.",
        "Minor enhancement (jam kerja terbatas).",
        "Response time 24 jam untuk issue blocker.",
        "Monitoring kesehatan sistem bulanan.",
    ])
    doc.add_page_break()

    # ============ 5. INVESTASI (HARGA DI AKHIR) ============
    H("5. INVESTASI & PEMBAYARAN", level=1)
    P("Setelah melihat output, outcome, dan potensi pendapatan Rp 102+ miliar, berikut rincian investasi yang dibutuhkan:", italic=True, color=GRAY)
    doc.add_paragraph()

    # Highlight contrast box
    t_contrast = doc.add_table(rows=2, cols=2); t_contrast.style = 'Light Grid Accent 1'
    t_contrast.cell(0,0).text = "Potensi PAD (skenario realistis)"
    t_contrast.cell(0,1).text = "Rp 102,7 miliar"
    t_contrast.cell(1,0).text = "Biaya investasi sistem"
    t_contrast.cell(1,1).text = rupiah(PRICE_TOTAL)
    set_cell_bg(t_contrast.cell(0,0), '28A745'); set_cell_bg(t_contrast.cell(0,1), '28A745')
    set_cell_bg(t_contrast.cell(1,0), '1C2D54'); set_cell_bg(t_contrast.cell(1,1), '1C2D54')
    for r in t_contrast.rows:
        for c in r.cells:
            for p in c.paragraphs:
                for run in p.runs: run.font.color.rgb = RGBColor(255,255,255); run.font.bold = True; run.font.size = Pt(13)
    doc.add_paragraph()
    P("Rasio ROI: >685× · Break-even: 15 PBG berhasil ditertibkan (0,015% dari 96.322 deteksi).", bold=True, color=SUCCESS, align='center', size=12)
    doc.add_paragraph()

    H("Rincian Harga per Paket", level=2)
    t = doc.add_table(rows=1, cols=3); t.style = 'Light Grid Accent 1'
    hdr = t.rows[0].cells
    hdr[0].text = "Paket"; hdr[1].text = "Deskripsi"; hdr[2].text = "Harga"
    for c in hdr:
        set_cell_bg(c, '1C2D54')
        for p in c.paragraphs:
            for r in p.runs: r.font.color.rgb = RGBColor(255,255,255); r.font.bold = True

    pkgs = [
        ("A", "Transfer, Deployment & Rekonsiliasi", rupiah(PRICE_A_TRANSFER)),
        ("B", "Pembangunan Fitur Monitoring Satelit", rupiah(PRICE_B_SATELIT)),
        ("C", "Data Enrichment Geocoding (opsional)", rupiah(PRICE_C_TRAINING)),
        ("D", "Maintenance & Support 6 Bulan (opsional)", rupiah(PRICE_D_MAINT)),
    ]
    for pkg in pkgs:
        row = t.add_row().cells
        for i, v in enumerate(pkg): row[i].text = v
    total_row = t.add_row().cells
    total_row[0].text = ""; total_row[1].text = "TOTAL (A+B+C+D)"; total_row[2].text = rupiah(PRICE_TOTAL)
    for c in total_row:
        set_cell_bg(c, 'C9A227')
        for p in c.paragraphs:
            for r in p.runs: r.font.bold = True

    doc.add_paragraph()
    H("Term Pembayaran", level=2)
    t = doc.add_table(rows=1, cols=3); t.style = 'Light Grid Accent 1'
    hdr = t.rows[0].cells
    hdr[0].text = "Milestone"; hdr[1].text = "%"; hdr[2].text = "Nominal"
    for c in hdr:
        set_cell_bg(c, '1C2D54')
        for p in c.paragraphs:
            for r in p.runs: r.font.color.rgb = RGBColor(255,255,255); r.font.bold = True
    for ms in [
        ("Down Payment (kickoff)", "30%", rupiah(int(PRICE_TOTAL*0.3))),
        ("Paket A selesai + data tersinkronisasi", "20%", rupiah(int(PRICE_TOTAL*0.2))),
        ("Paket B backend + peta live", "30%", rupiah(int(PRICE_TOTAL*0.3))),
        ("Handover + UAT sign-off", "20%", rupiah(int(PRICE_TOTAL*0.2))),
    ]:
        row = t.add_row().cells
        for i, v in enumerate(ms): row[i].text = v

    doc.add_paragraph()
    P("Catatan:", bold=True)
    for b in [
        "Harga sudah termasuk source code + dokumentasi lengkap.",
        "Deployment di infrastruktur klien (tidak include biaya server/domain).",
        "Garansi bug-fix 60 hari setelah handover.",
        "Perubahan scope di luar proposal akan di-quote terpisah sebagai change-request.",
        "Penawaran berlaku 30 hari sejak tanggal diterbitkan.",
    ]:
        doc.add_paragraph(b, style='List Bullet')
    doc.add_page_break()

    # ============ 6. TIMELINE ============
    H("6. TIMELINE", level=1)
    t = doc.add_table(rows=1, cols=3); t.style = 'Light Grid Accent 1'
    hdr = t.rows[0].cells
    hdr[0].text = "Fase"; hdr[1].text = "Minggu"; hdr[2].text = "Deliverable"
    for c in hdr:
        set_cell_bg(c, '1C2D54')
        for p in c.paragraphs:
            for r in p.runs: r.font.color.rgb = RGBColor(255,255,255); r.font.bold = True
    for ph in [
        ("Transfer & Deploy", "1-2", "Production live, data tersinkronisasi."),
        ("Data Pipeline & Integration", "3-4", "1.09M row terklasifikasi ke 31 kecamatan."),
        ("Backend API", "5-6", "7 endpoint siap, cache strategy aktif."),
        ("Frontend Map & Dashboard", "7-9", "Peta interaktif + filter + dashboard live."),
        ("Verification & Testing", "10", "Workflow verify + SQL audit proof."),
        ("UAT + Handover", "11", "Training + dokumentasi + sign-off."),
    ]:
        row = t.add_row().cells
        for i, v in enumerate(ph): row[i].text = v

    doc.add_paragraph()
    P("Total durasi: 11 minggu (~2,5 bulan).", bold=True, color=NAVY)
    doc.add_page_break()

    # ============ PENUTUP ============
    H("7. PENUTUP", level=1)
    P("Proposal ini menawarkan solusi komprehensif untuk monitoring bangunan di Kabupaten Bandung dengan potensi ROI sangat tinggi — minimal 685× pada skenario realistis. Dari data real SIMBG, fitur ini balik modal cukup dengan 15 PBG berhasil ditertibkan.")
    doc.add_paragraph()
    P("Kami berkomitmen menghasilkan sistem yang production-ready, proven dengan proof-of-correctness SQL audit, dan didukung dokumentasi lengkap untuk kelanjutan pemeliharaan internal DPUTR.")
    doc.add_paragraph(); doc.add_paragraph()
    P("Terima kasih atas kesempatan ini.", align='center', italic=True)
    doc.add_paragraph(); doc.add_paragraph()

    t = doc.add_table(rows=2, cols=2); t.autofit = False
    t.columns[0].width = Cm(7); t.columns[1].width = Cm(7)
    t.cell(0,0).text = "Hormat kami,"; t.cell(0,1).text = f"{CLIENT_NAME},"
    t.cell(1,0).text = f"\n\n\n{VENDOR_NAME}"
    t.cell(1,1).text = "\n\n\n__________________________"

    out = os.path.join(OUT_DIR, 'Proposal_Monitoring_Satelit.docx')
    doc.save(out)
    print(f"✓ DOCX: {out}")
    return out


# ====================== PDF ======================
def build_pdf():
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm, mm
    from reportlab.lib.colors import HexColor, white, black
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, KeepTogether
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY

    NAVY_HEX = HexColor('#1C2D54'); GOLD_HEX = HexColor('#C9A227'); LIGHT_HEX = HexColor('#E8ECEF')
    GRAY_HEX = HexColor('#5A5A5A'); SUCCESS_HEX = HexColor('#28A745')

    out = os.path.join(OUT_DIR, 'Proposal_Monitoring_Satelit.pdf')
    doc = SimpleDocTemplate(out, pagesize=A4, leftMargin=2.5*cm, rightMargin=2*cm,
                            topMargin=2.2*cm, bottomMargin=2*cm)
    styles = getSampleStyleSheet()
    h1 = ParagraphStyle('h1', parent=styles['Heading1'], fontSize=22, textColor=NAVY_HEX, spaceAfter=14, leading=26, fontName='Helvetica-Bold')
    h2 = ParagraphStyle('h2', parent=styles['Heading2'], fontSize=15, textColor=NAVY_HEX, spaceAfter=8, fontName='Helvetica-Bold')
    h3 = ParagraphStyle('h3', parent=styles['Heading3'], fontSize=12, textColor=NAVY_HEX, spaceAfter=6, fontName='Helvetica-Bold')
    body = ParagraphStyle('body', parent=styles['Normal'], fontSize=11, leading=15, spaceAfter=6, fontName='Helvetica', alignment=TA_JUSTIFY)
    bullet = ParagraphStyle('bullet', parent=body, leftIndent=18, bulletIndent=6, spaceAfter=3)
    center = ParagraphStyle('center', parent=body, alignment=TA_CENTER)
    big_center = ParagraphStyle('big_center', parent=body, alignment=TA_CENTER, fontSize=16, leading=20, textColor=NAVY_HEX)
    cover_title = ParagraphStyle('ct', fontSize=32, alignment=TA_CENTER, textColor=NAVY_HEX, fontName='Helvetica-Bold', leading=40, spaceAfter=10)
    cover_sub = ParagraphStyle('cs', fontSize=18, alignment=TA_CENTER, textColor=GOLD_HEX, fontName='Helvetica-Bold', leading=24, spaceAfter=16)
    cover_tag = ParagraphStyle('ctag', fontSize=13, alignment=TA_CENTER, textColor=NAVY_HEX, fontName='Helvetica-Oblique', leading=18)
    quote_style = ParagraphStyle('q', parent=body, alignment=TA_CENTER, fontSize=12, textColor=NAVY_HEX, fontName='Helvetica-Oblique')

    def rupiah(n): return f"Rp {n:,.0f}".replace(",", ".")

    story = []

    # COVER
    story.append(Spacer(1, 4*cm))
    story.append(Paragraph("PROPOSAL PENAWARAN", cover_title))
    story.append(Paragraph(PROJECT_NAME, cover_sub))
    story.append(Spacer(1, 1*cm))
    story.append(Paragraph("Diajukan kepada:", cover_tag))
    story.append(Spacer(1, 4*mm))
    story.append(Paragraph(f"<b>{CLIENT_NAME}</b>", ParagraphStyle('cn', fontSize=16, alignment=TA_CENTER, textColor=NAVY_HEX, fontName='Helvetica-Bold')))
    story.append(Spacer(1, 2*cm))

    meta = [
        ["Nomor Proposal", PROPOSAL_NO],
        ["Tanggal", DATE_STR],
        ["Penyedia", VENDOR_NAME],
    ]
    t = Table(meta, colWidths=[5*cm, 10*cm])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (0,-1), LIGHT_HEX),
        ('BOX', (0,0), (-1,-1), 0.5, GRAY_HEX),
        ('INNERGRID', (0,0), (-1,-1), 0.25, HexColor('#BFC7D0')),
        ('FONTNAME', (0,0), (0,-1), 'Helvetica-Bold'),
        ('FONTSIZE', (0,0), (-1,-1), 11),
        ('LEFTPADDING', (0,0), (-1,-1), 8),
        ('TOPPADDING', (0,0), (-1,-1), 6),
        ('BOTTOMPADDING', (0,0), (-1,-1), 6),
        ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
    ]))
    story.append(t)
    story.append(Spacer(1, 2*cm))
    story.append(Paragraph('"Fitur ini membayar dirinya sendiri kalau DPUTR berhasil', quote_style))
    story.append(Paragraph('menertibkan 15 bangunan saja dari 96.322 yang terdeteksi."', quote_style))
    story.append(PageBreak())

    # EXECUTIVE SUMMARY
    story.append(Paragraph("RINGKASAN EKSEKUTIF", h1))
    story.append(Paragraph(
        "Kabupaten Bandung memiliki ribuan bangunan yang belum memiliki Persetujuan Bangunan Gedung (PBG) — "
        "sumber potensi pendapatan asli daerah (PAD) yang selama ini tidak termonitor. Sistem Monitoring Satelit "
        "ini memanfaatkan citra satelit + data SIMBG untuk mendeteksi, mengklasifikasi, dan memetakan seluruh "
        "bangunan di 31 kecamatan Kabupaten Bandung secara otomatis.", body))
    story.append(Spacer(1, 6*mm))
    story.append(Paragraph("Angka Kunci", h2))
    kv_data = [
        ["Total bangunan terdeteksi", "1,09 juta"],
        ["Bangunan ≥200 m² tanpa izin sah", "96.322"],
        ["Potensi PAD (skenario 10%)", "Rp 102,7 miliar"],
    ]
    t = Table(kv_data, colWidths=[8*cm, 7*cm])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (0,-1), LIGHT_HEX),
        ('BOX', (0,0), (-1,-1), 0.75, NAVY_HEX),
        ('INNERGRID', (0,0), (-1,-1), 0.25, GRAY_HEX),
        ('FONTNAME', (0,0), (0,-1), 'Helvetica-Bold'),
        ('FONTSIZE', (0,0), (-1,-1), 11),
        ('LEFTPADDING', (0,0), (-1,-1), 8),
        ('TOPPADDING', (0,0), (-1,-1), 8),
        ('BOTTOMPADDING', (0,0), (-1,-1), 8),
    ]))
    story.append(t)
    story.append(Spacer(1, 6*mm))
    story.append(Paragraph('<i>Dari 5.241 PBG yang sudah terbit di Kab. Bandung, rata-rata retribusi yang masuk = Rp 10,66 juta per PBG.</i>', body))
    story.append(Paragraph('<b><font color="#28A745">Artinya: 96.322 bangunan yang belum tertangani setara dengan Rp 1,03 triliun potensi PAD yang belum dieksekusi.</font></b>', body))
    story.append(PageBreak())

    # LATAR BELAKANG
    story.append(Paragraph("1. LATAR BELAKANG", h1))
    story.append(Paragraph("DPUTR Kabupaten Bandung membutuhkan instrumen monitoring yang dapat:", body))
    for b in [
        "Mendeteksi bangunan tanpa izin di 31 kecamatan secara otomatis.",
        "Menyajikan data visual peta interaktif untuk staff enforcement.",
        "Mencocokkan data deteksi dengan record PBG SIMBG existing.",
        "Memberikan dasar kuantitatif untuk target PAD tahunan.",
    ]:
        story.append(Paragraph(f"• {b}", bullet))
    story.append(Spacer(1, 6*mm))
    story.append(Paragraph(
        "Metode konvensional (survey lapangan) memakan waktu berbulan-bulan dan biaya tinggi. Solusi berbasis "
        "satelit + SIMBG menghasilkan data yang sama dalam hitungan menit dan bisa di-refresh kapan saja.", body))
    story.append(PageBreak())

    # ============ 2. OUTPUT & OUTCOME ============
    story.append(Paragraph("2. OUTPUT &amp; OUTCOME", h1))
    story.append(Paragraph("Output — Yang DPUTR Terima", h2))
    for b in [
        "Aplikasi web Monitoring Satelit terdeploy di production.",
        "Peta interaktif 31 polygon + 97.000+ titik deteksi bangunan.",
        "Dashboard KPI realtime + tabel breakdown per kecamatan.",
        "Workflow verifikasi dengan audit trail.",
        "Source code + database schema + dokumentasi lengkap.",
        "1 batch training untuk staff DPUTR.",
    ]:
        story.append(Paragraph(f"• {b}", bullet))
    story.append(Spacer(1, 4*mm))
    story.append(Paragraph("Outcome — Manfaat Terukur", h2))
    for b in [
        "Visibilitas 96.322 bangunan potensial wajib PBG yang selama ini tidak termonitor.",
        "Efisiensi enforcement: target langsung dipetakan, tidak perlu survey acak.",
        "Basis data terpadu untuk sinkronisasi dengan PBB, tata ruang, dan perizinan.",
        "Data-driven decision untuk kebijakan PAD.",
        "Dokumentasi kontribusi kinerja DPUTR dalam laporan tahunan.",
    ]:
        story.append(Paragraph(f"• {b}", bullet))
    story.append(PageBreak())

    # ============ 3. POTENSI PENDAPATAN ============
    story.append(Paragraph("3. POTENSI PENDAPATAN (ROI)", h1))
    story.append(Paragraph("<i>Berdasarkan data real SIMBG Kab. Bandung (5.241 PBG yang sudah terbit):</i>", body))
    fact_data = [
        ["Rata-rata retribusi per PBG", "Rp 10.660.000"],
        ["Total PAD PBG terkumpul di DB", "Rp 55,86 miliar"],
    ]
    t = Table(fact_data, colWidths=[8*cm, 7*cm])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (0,-1), LIGHT_HEX),
        ('BOX', (0,0), (-1,-1), 0.75, NAVY_HEX),
        ('INNERGRID', (0,0), (-1,-1), 0.25, GRAY_HEX),
        ('FONTNAME', (0,0), (0,-1), 'Helvetica-Bold'),
        ('FONTSIZE', (0,0), (-1,-1), 11),
        ('LEFTPADDING', (0,0), (-1,-1), 8),
        ('TOPPADDING', (0,0), (-1,-1), 8),
        ('BOTTOMPADDING', (0,0), (-1,-1), 8),
    ]))
    story.append(t)
    story.append(Spacer(1, 6*mm))

    story.append(Paragraph("Skenario Potensi Pendapatan dari 96.322 Bangunan Tanpa Izin", h2))
    scen_data = [
        ["Skenario Enforcement", "PBG Tertarik", "Potensi PAD"],
        ["1% (super konservatif)", "963", "Rp 10,3 miliar"],
        ["5% (konservatif)", "4.816", "Rp 51,3 miliar"],
        ["10% (realistis)", "9.632", "Rp 102,7 miliar"],
        ["20% (target wajar)", "19.264", "Rp 205,4 miliar"],
        ["30% (agresif)", "28.897", "Rp 308,0 miliar"],
    ]
    t = Table(scen_data, colWidths=[6*cm, 4*cm, 5*cm])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), NAVY_HEX),
        ('TEXTCOLOR', (0,0), (-1,0), white),
        ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
        ('FONTSIZE', (0,0), (-1,-1), 10.5),
        ('BOX', (0,0), (-1,-1), 0.75, NAVY_HEX),
        ('INNERGRID', (0,0), (-1,-1), 0.25, GRAY_HEX),
        ('LEFTPADDING', (0,0), (-1,-1), 6),
        ('TOPPADDING', (0,0), (-1,-1), 6),
        ('BOTTOMPADDING', (0,0), (-1,-1), 6),
        ('ROWBACKGROUNDS', (0,1), (-1,-1), [white, HexColor('#F5F7FA')]),
        ('ALIGN', (1,1), (-1,-1), 'RIGHT'),
    ]))
    story.append(t)
    story.append(Spacer(1, 6*mm))
    story.append(Paragraph('<i><font color="#1C2D54">Skenario realistis saja (10%) sudah setara membangun 3 jembatan baru.</font></i>', center))
    story.append(Spacer(1, 6*mm))
    story.append(Paragraph("Manfaat Tambahan (Belum Dihitung)", h2))
    for b in [
        "Denda administratif bangunan tanpa izin: 2-10% nilai bangunan.",
        "Akurasi penagihan PBB: potensi tambahan Rp 10-20 miliar/tahun.",
        "Compliance tata ruang & RDTR enforcement.",
        "Decision support untuk perencanaan infrastruktur.",
    ]:
        story.append(Paragraph(f"• {b}", bullet))
    story.append(PageBreak())

    # ============ 4. RUANG LINGKUP ============
    story.append(Paragraph("4. RUANG LINGKUP PEKERJAAN", h1))
    story.append(Paragraph("Seluruh scope pekerjaan disajikan dalam bentuk tabel untuk kemudahan audit dan tracking deliverable. Rincian nilai investasi ada pada Bagian 5.", body))
    story.append(Spacer(1, 4*mm))

    def pkg_as_table(code, name, items):
        """Render paket sederhana sebagai tabel TANPA harga."""
        data = [
            ["Kode Paket", code],
            ["Nama Paket", name],
            ["Deliverable & Fitur", ""],
        ]
        for i, it in enumerate(items):
            data.append([f"{code}.{i+1}", it])
        t = Table(data, colWidths=[3.5*cm, 12*cm])
        style = [
            ('BACKGROUND', (0,0), (0,1), NAVY_HEX),
            ('TEXTCOLOR', (0,0), (0,1), white),
            ('FONTNAME', (0,0), (0,1), 'Helvetica-Bold'),
            ('BACKGROUND', (0,2), (-1,2), GOLD_HEX),
            ('TEXTCOLOR', (0,2), (-1,2), white),
            ('FONTNAME', (0,2), (-1,2), 'Helvetica-Bold'),
            ('SPAN', (0,2), (-1,2)),
            ('BACKGROUND', (0,3), (0,-1), LIGHT_HEX),
            ('FONTNAME', (0,3), (0,-1), 'Helvetica-Bold'),
            ('TEXTCOLOR', (0,3), (0,-1), NAVY_HEX),
            ('BOX', (0,0), (-1,-1), 0.75, NAVY_HEX),
            ('INNERGRID', (0,0), (-1,-1), 0.25, GRAY_HEX),
            ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
            ('ALIGN', (0,0), (0,-1), 'CENTER'),
            ('LEFTPADDING', (0,0), (-1,-1), 6),
            ('RIGHTPADDING', (0,0), (-1,-1), 6),
            ('TOPPADDING', (0,0), (-1,-1), 5),
            ('BOTTOMPADDING', (0,0), (-1,-1), 5),
            ('FONTSIZE', (0,0), (-1,-1), 10.5),
        ]
        t.setStyle(TableStyle(style))
        story.append(t)
        story.append(Spacer(1, 6*mm))

    # Paket A
    story.append(Paragraph("Paket A — Transfer, Deployment & Rekonsiliasi", h2))
    pkg_as_table("A", "Transfer, Deployment & Rekonsiliasi Data", [
        "Setup production environment (nginx, PHP-FPM, MariaDB tuning).",
        "Import & migrate database produksi (297 MB).",
        "Bug-fix login session, Vite pipeline, asset loading.",
        "Rekonsiliasi angka dengan Dashboard Pimpinan SIMBG & Dashboard PBG.",
        "Deployment runbook + rollback procedure.",
        "Handover documentation + 1 batch training.",
    ])

    # Paket B — big table dgn sub-paket
    story.append(Paragraph("Paket B — Pembangunan Fitur Monitoring Satelit", h2))
    story.append(Paragraph("Dipecah menjadi 7 sub-paket teknis:", body))
    story.append(Spacer(1, 3*mm))

    subpkg = [
        ("B.1", "Data Pipeline & Integrasi", [
            "Ingest dataset Microsoft Building Footprints (1,09 juta bangunan).",
            "Integrasi data PBG SIMBG (pbg_task, pbg_task_details, retributions).",
            "Spatial matching engine deteksi satelit ↔ PBG.",
            "Integrasi polygon resmi 31 kecamatan Kab. Bandung (sumber BPS).",
            "Point-in-polygon classification 1,09M rows via Python shapely.",
            "Data quality audit (orphan FK, validasi spasial).",
        ]),
        ("B.2", "Backend API", [
            "7 REST endpoint: stats, geojson, pbg-geojson, verify, refresh.",
            "Multi-filter query (kecamatan, jenis, status, luas).",
            "2-tier caching (Laravel cache + DB snapshot).",
            "Indexing & query optimization (cold 15s → warm 0,4s).",
            "Authentication layer (Sanctum + role-based).",
        ]),
        ("B.3", "Frontend Peta Interaktif", [
            "Map engine Leaflet + Esri satellite tile + Carto labels.",
            "Bounded scroll + dynamic min-zoom (area lock Kab. Bandung).",
            "Render 31 polygon kecamatan + click-to-zoom + tooltip.",
            "Marker clustering 2-layer (deteksi satelit + titik PBG).",
            "Popup detail bangunan + panel verifikasi.",
            "Loading indicator + abort-on-pan race handling.",
            "Legend 5-state (SK Terbit / Proses / Ditolak / Tanpa Izin / PBG).",
            "Filter system 6-kategori (kecamatan, jenis, status, luas, tampilan, layer).",
        ]),
        ("B.4", "Dashboard Analytics", [
            "4 KPI cards realtime (Total / Berizin / Tanpa Izin / Rasio).",
            "PBG breakdown strip (Total / Terbit / Proses / Ditolak).",
            "Tabel 'Data per Kecamatan' dengan stacked bar 3-segment.",
            "Grand total row + snapshot timestamp.",
            "Breakdown jenis bangunan (7 kategori fungsi).",
            "Drill-down interaktif (klik row → zoom + filter).",
        ]),
        ("B.5", "Verification & Audit Workflow", [
            "Panel verifikasi 4-status (Illegal / Legal / Review / False+).",
            "Bulk verification API + audit trail (verified_by, verified_at, notes).",
            "Role-based access control.",
        ]),
        ("B.6", "Performance & Production Readiness", [
            "Tabel snapshot kecamatan_stats (pre-computed, 186 rows).",
            "Artisan command kecamatan-stats:refresh (idempotent, ~60s).",
            "Cache strategy dengan TTL + auto-invalidation on verify.",
        ]),
        ("B.7", "Dokumentasi & Proof-of-Correctness", [
            "SQL audit script verifikasi angka.",
            "User manual staff PUTR (screenshot + workflow).",
            "Technical documentation (schema ERD, API spec).",
        ]),
    ]

    # Build consolidated Paket B table
    b_data = [["Kode", "Sub-Paket / Deliverable"]]
    b_style_spans = []  # (row_idx)
    gold_rows = []  # section header rows (sub-paket title)
    for code, name, items in subpkg:
        gold_rows.append(len(b_data))
        b_data.append([code, name])
        for j, it in enumerate(items):
            b_data.append([f"{code}.{j+1}", it])

    t = Table(b_data, colWidths=[2.5*cm, 13*cm])
    style = [
        ('BACKGROUND', (0,0), (-1,0), NAVY_HEX),
        ('TEXTCOLOR', (0,0), (-1,0), white),
        ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
        ('BOX', (0,0), (-1,-1), 0.75, NAVY_HEX),
        ('INNERGRID', (0,0), (-1,-1), 0.25, GRAY_HEX),
        ('FONTSIZE', (0,0), (-1,-1), 10),
        ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
        ('ALIGN', (0,0), (0,-1), 'CENTER'),
        ('LEFTPADDING', (0,0), (-1,-1), 5),
        ('RIGHTPADDING', (0,0), (-1,-1), 5),
        ('TOPPADDING', (0,0), (-1,-1), 5),
        ('BOTTOMPADDING', (0,0), (-1,-1), 5),
    ]
    # Gold backgrounds pada sub-paket header rows
    for gr in gold_rows:
        style.append(('BACKGROUND', (0, gr), (-1, gr), GOLD_HEX))
        style.append(('TEXTCOLOR', (0, gr), (-1, gr), white))
        style.append(('FONTNAME', (0, gr), (-1, gr), 'Helvetica-Bold'))
        style.append(('FONTSIZE', (0, gr), (-1, gr), 11))
    # Deliverable-only rows (bukan header & bukan section): light bg kolom-1
    deliv_rows = [i for i in range(1, len(b_data)) if i not in gold_rows]
    for dr in deliv_rows:
        style.append(('BACKGROUND', (0, dr), (0, dr), LIGHT_HEX))
        style.append(('FONTNAME', (0, dr), (0, dr), 'Helvetica-Bold'))
        style.append(('TEXTCOLOR', (0, dr), (0, dr), NAVY_HEX))
    t.setStyle(TableStyle(style))
    story.append(t)
    story.append(Spacer(1, 6*mm))

    # Paket C & D (no price)
    story.append(Paragraph("Paket C — Data Enrichment Geocoding (Opsional)", h2))
    pkg_as_table("C", "Enrichment 96.322 Bangunan dengan Alamat Lengkap", [
        "Reverse geocoding seluruh koordinat deteksi ke alamat (jalan, RT/RW, kelurahan).",
        "Integrasi Badan Informasi Geospasial (BIG) Tanah Air sebagai source utama (gratis).",
        "Fallback Nominatim OpenStreetMap untuk area yang tidak tercover BIG.",
        "Batch processing 96k+ row dengan queue + retry handling.",
        "Cache hasil geocoding di DB untuk efisiensi query lanjutan.",
        "Tampil di popup marker + export CSV lengkap alamat.",
        "Export ready-to-use untuk surat teguran (template mail-merge).",
    ])

    story.append(Paragraph("Paket D — Maintenance & Support 6 Bulan (Opsional)", h2))
    pkg_as_table("D", "Maintenance & Support 6 Bulan", [
        "Bug-fix unlimited untuk defect.",
        "Minor enhancement (jam kerja terbatas).",
        "Response time 24 jam untuk issue blocker.",
        "Monitoring kesehatan sistem bulanan.",
    ])
    story.append(PageBreak())

    # ============ 5. INVESTASI & PEMBAYARAN ============
    story.append(Paragraph("5. INVESTASI &amp; PEMBAYARAN", h1))
    story.append(Paragraph("<i>Setelah melihat output, outcome, dan potensi pendapatan Rp 102+ miliar, berikut rincian investasi yang dibutuhkan:</i>", body))
    story.append(Spacer(1, 4*mm))

    # Contrast highlight
    contrast_data = [
        ["Potensi PAD (skenario realistis)", "Rp 102,7 miliar"],
        ["Biaya investasi sistem", rupiah(PRICE_TOTAL)],
    ]
    t = Table(contrast_data, colWidths=[8*cm, 7*cm])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), HexColor('#28A745')),
        ('BACKGROUND', (0,1), (-1,1), NAVY_HEX),
        ('TEXTCOLOR', (0,0), (-1,-1), white),
        ('FONTNAME', (0,0), (-1,-1), 'Helvetica-Bold'),
        ('FONTSIZE', (0,0), (-1,-1), 13),
        ('BOX', (0,0), (-1,-1), 0.75, NAVY_HEX),
        ('INNERGRID', (0,0), (-1,-1), 0.25, white),
        ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
        ('LEFTPADDING', (0,0), (-1,-1), 10),
        ('TOPPADDING', (0,0), (-1,-1), 10),
        ('BOTTOMPADDING', (0,0), (-1,-1), 10),
    ]))
    story.append(t)
    story.append(Spacer(1, 6*mm))
    story.append(Paragraph('<b><font color="#28A745">Rasio ROI: &gt;685× · Break-even: 15 PBG berhasil ditertibkan (0,015% dari 96.322 deteksi).</font></b>', center))
    story.append(Spacer(1, 6*mm))

    story.append(Paragraph("Rincian Harga per Paket", h2))
    price_data = [
        ["Paket", "Deskripsi", "Harga"],
        ["A", "Transfer, Deployment & Rekonsiliasi", rupiah(PRICE_A_TRANSFER)],
        ["B", "Pembangunan Fitur Monitoring Satelit", rupiah(PRICE_B_SATELIT)],
        ["C", "Data Enrichment Geocoding (opsional)", rupiah(PRICE_C_TRAINING)],
        ["D", "Maintenance & Support 6 Bulan (opsional)", rupiah(PRICE_D_MAINT)],
        ["", "TOTAL (A+B+C+D)", rupiah(PRICE_TOTAL)],
    ]
    t = Table(price_data, colWidths=[2*cm, 9*cm, 4*cm])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), NAVY_HEX),
        ('TEXTCOLOR', (0,0), (-1,0), white),
        ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
        ('BACKGROUND', (0,-1), (-1,-1), GOLD_HEX),
        ('FONTNAME', (0,-1), (-1,-1), 'Helvetica-Bold'),
        ('FONTSIZE', (0,0), (-1,-1), 11),
        ('BOX', (0,0), (-1,-1), 0.75, NAVY_HEX),
        ('INNERGRID', (0,0), (-1,-1), 0.25, GRAY_HEX),
        ('ALIGN', (0,0), (0,-1), 'CENTER'),
        ('ALIGN', (2,1), (2,-1), 'RIGHT'),
        ('LEFTPADDING', (0,0), (-1,-1), 6),
        ('TOPPADDING', (0,0), (-1,-1), 7),
        ('BOTTOMPADDING', (0,0), (-1,-1), 7),
    ]))
    story.append(t)
    story.append(Spacer(1, 8*mm))

    story.append(Paragraph("Term Pembayaran", h2))
    pay_data = [
        ["Milestone", "%", "Nominal"],
        ["Down Payment (kickoff)", "30%", rupiah(int(PRICE_TOTAL*0.3))],
        ["Paket A selesai + data tersinkronisasi", "20%", rupiah(int(PRICE_TOTAL*0.2))],
        ["Paket B backend + peta live", "30%", rupiah(int(PRICE_TOTAL*0.3))],
        ["Handover + UAT sign-off", "20%", rupiah(int(PRICE_TOTAL*0.2))],
    ]
    t = Table(pay_data, colWidths=[9*cm, 2*cm, 4*cm])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), NAVY_HEX),
        ('TEXTCOLOR', (0,0), (-1,0), white),
        ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
        ('FONTSIZE', (0,0), (-1,-1), 11),
        ('BOX', (0,0), (-1,-1), 0.75, NAVY_HEX),
        ('INNERGRID', (0,0), (-1,-1), 0.25, GRAY_HEX),
        ('ROWBACKGROUNDS', (0,1), (-1,-1), [white, HexColor('#F5F7FA')]),
        ('ALIGN', (1,0), (-1,-1), 'CENTER'),
        ('ALIGN', (2,1), (2,-1), 'RIGHT'),
    ]))
    story.append(t)
    story.append(Spacer(1, 6*mm))

    story.append(Paragraph("Catatan:", h3))
    for b in [
        "Harga sudah termasuk source code + dokumentasi lengkap.",
        "Deployment di infrastruktur klien (tidak include biaya server/domain).",
        "Garansi bug-fix 60 hari setelah handover.",
        "Perubahan scope di luar proposal akan di-quote terpisah.",
        "Penawaran berlaku 30 hari sejak tanggal diterbitkan.",
    ]:
        story.append(Paragraph(f"• {b}", bullet))
    story.append(PageBreak())

    # TIMELINE
    story.append(Paragraph("6. TIMELINE", h1))
    timeline_data = [
        ["Fase", "Minggu", "Deliverable"],
        ["Transfer & Deploy", "1-2", "Production live, data tersinkronisasi."],
        ["Data Pipeline & Integration", "3-4", "1,09M row terklasifikasi ke 31 kecamatan."],
        ["Backend API", "5-6", "7 endpoint siap, cache strategy aktif."],
        ["Frontend Map & Dashboard", "7-9", "Peta + filter + dashboard live."],
        ["Verification & Testing", "10", "Workflow verify + SQL audit proof."],
        ["UAT + Handover", "11", "Training + dokumentasi + sign-off."],
    ]
    t = Table(timeline_data, colWidths=[5*cm, 2.5*cm, 7.5*cm])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), NAVY_HEX),
        ('TEXTCOLOR', (0,0), (-1,0), white),
        ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
        ('FONTSIZE', (0,0), (-1,-1), 10.5),
        ('BOX', (0,0), (-1,-1), 0.75, NAVY_HEX),
        ('INNERGRID', (0,0), (-1,-1), 0.25, GRAY_HEX),
        ('ROWBACKGROUNDS', (0,1), (-1,-1), [white, HexColor('#F5F7FA')]),
        ('ALIGN', (1,1), (1,-1), 'CENTER'),
        ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
        ('LEFTPADDING', (0,0), (-1,-1), 6),
        ('TOPPADDING', (0,0), (-1,-1), 7),
        ('BOTTOMPADDING', (0,0), (-1,-1), 7),
    ]))
    story.append(t)
    story.append(Spacer(1, 6*mm))
    story.append(Paragraph('<b><font color="#1C2D54">Total durasi: 11 minggu (~2,5 bulan).</font></b>', body))
    story.append(PageBreak())

    # PENUTUP
    story.append(Paragraph("7. PENUTUP", h1))
    story.append(Paragraph(
        "Proposal ini menawarkan solusi komprehensif untuk monitoring bangunan di Kabupaten Bandung dengan potensi "
        "ROI sangat tinggi — minimal 685x pada skenario realistis. Dari data real SIMBG, fitur ini balik modal cukup "
        "dengan 15 PBG berhasil ditertibkan.", body))
    story.append(Spacer(1, 4*mm))
    story.append(Paragraph(
        "Kami berkomitmen menghasilkan sistem yang production-ready, proven dengan proof-of-correctness SQL audit, "
        "dan didukung dokumentasi lengkap untuk kelanjutan pemeliharaan internal DPUTR.", body))
    story.append(Spacer(1, 8*mm))
    story.append(Paragraph("<i>Terima kasih atas kesempatan ini.</i>", center))
    story.append(Spacer(1, 2*cm))

    sig = [
        ["Hormat kami,", f"{CLIENT_NAME},"],
        ["", ""],
        ["", ""],
        ["", ""],
        [VENDOR_NAME, "__________________________"],
    ]
    t = Table(sig, colWidths=[8*cm, 8*cm])
    t.setStyle(TableStyle([
        ('FONTSIZE', (0,0), (-1,-1), 11),
        ('BOTTOMPADDING', (0,0), (-1,-1), 8),
    ]))
    story.append(t)

    doc.build(story)
    print(f"✓ PDF:  {out}")
    return out


# ====================== PPTX ======================
def build_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

    NAVY_RGB = RGBColor(28, 45, 84); GOLD_RGB = RGBColor(201, 162, 39)
    WHITE_RGB = RGBColor(255, 255, 255); GRAY_RGB = RGBColor(90, 90, 90)
    SUCCESS_RGB = RGBColor(40, 167, 69); LIGHT_RGB = RGBColor(232, 236, 239)
    DANGER_RGB = RGBColor(220, 53, 69)

    blank_layout = prs.slide_layouts[6]

    def add_bg(slide, color=WHITE_RGB):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid(); bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        return bg

    def add_bar(slide, x, y, w, h, color):
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        bar.fill.solid(); bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        return bar

    def add_text(slide, text, x, y, w, h, size=18, bold=False, color=NAVY_RGB, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font='Calibri'):
        tx = slide.shapes.add_textbox(x, y, w, h)
        tf = tx.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
        tf.margin_left = Emu(0); tf.margin_right = Emu(0); tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = font
        return tx

    def add_multiline(slide, lines, x, y, w, h, size=14, bold=False, color=NAVY_RGB, align=PP_ALIGN.LEFT, font='Calibri'):
        tx = slide.shapes.add_textbox(x, y, w, h)
        tf = tx.text_frame; tf.word_wrap = True
        tf.margin_left = Emu(0); tf.margin_right = Emu(0); tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            r = p.add_run(); r.text = line
            r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = font
        return tx

    def add_header(slide, title, num=None):
        add_bar(slide, 0, 0, prs.slide_width, Inches(0.8), NAVY_RGB)
        if num:
            add_text(slide, num, Inches(0.5), Inches(0.2), Inches(1), Inches(0.4), size=14, bold=True, color=GOLD_RGB)
        add_text(slide, title, Inches(1.4) if num else Inches(0.5), Inches(0.15), Inches(11), Inches(0.5), size=22, bold=True, color=WHITE_RGB)
        add_bar(slide, 0, Inches(0.8), prs.slide_width, Inches(0.05), GOLD_RGB)

    # ============ SLIDE 1: COVER ============
    s = prs.slides.add_slide(blank_layout)
    add_bg(s, NAVY_RGB)
    # Gold accent bar
    add_bar(s, 0, Inches(3.0), prs.slide_width, Inches(0.1), GOLD_RGB)
    add_text(s, "PROPOSAL PENAWARAN", Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.8), size=28, bold=True, color=WHITE_RGB, align=PP_ALIGN.CENTER)
    add_text(s, PROJECT_NAME, Inches(0.5), Inches(2.6), Inches(12.3), Inches(0.6), size=20, color=GOLD_RGB, align=PP_ALIGN.CENTER, bold=True)
    add_text(s, f"Untuk {CLIENT_NAME}", Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.5), size=18, color=WHITE_RGB, align=PP_ALIGN.CENTER)
    add_text(s, DATE_STR, Inches(0.5), Inches(4.8), Inches(12.3), Inches(0.4), size=14, color=LIGHT_RGB, align=PP_ALIGN.CENTER)
    add_text(s, f"Disiapkan oleh: {VENDOR_NAME}", Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.3), size=11, color=LIGHT_RGB, align=PP_ALIGN.CENTER)

    # ============ SLIDE 2: MASALAH ============
    s = prs.slides.add_slide(blank_layout); add_bg(s)
    add_header(s, "Kenapa Kami Butuh Sistem Ini?", "01")
    add_text(s, "Situasi Saat Ini", Inches(0.5), Inches(1.2), Inches(12), Inches(0.5), size=20, bold=True)
    problems = [
        ("96.322", "bangunan ≥200 m² di Kab. Bandung tidak tercatat PBG-nya", DANGER_RGB),
        ("Manual", "enforcement lewat survey lapangan: lambat, mahal, bias", GRAY_RGB),
        ("Tidak ada", "monitoring sistematis per kecamatan — sulit prioritisasi", GRAY_RGB),
        ("Potensi PAD", "miliaran rupiah hilang tiap tahun", GOLD_RGB),
    ]
    for i, (num, txt, color) in enumerate(problems):
        y = Inches(2.0 + i * 1.1)
        # big number
        add_text(s, num, Inches(0.7), y, Inches(3), Inches(0.8), size=26, bold=True, color=color, align=PP_ALIGN.LEFT)
        # desc
        add_text(s, txt, Inches(4.0), y + Inches(0.15), Inches(8.5), Inches(0.6), size=16, color=NAVY_RGB)

    # ============ SLIDE 3: SOLUSI ============
    s = prs.slides.add_slide(blank_layout); add_bg(s)
    add_header(s, "Solusi: Sistem Monitoring Satelit", "02")
    add_text(s, "Deteksi + Klasifikasi + Verifikasi — semua dalam satu dashboard", Inches(0.5), Inches(1.1), Inches(12), Inches(0.5), size=15, color=GRAY_RGB)

    steps = [
        ("🛰️", "DETEKSI", "1,09 juta bangunan dari citra satelit Microsoft Footprints"),
        ("📍", "KLASIFIKASI", "Point-in-polygon BPS → 31 kecamatan"),
        ("🔗", "MATCHING", "Cocokkan dengan PBG SIMBG → status izin tiap bangunan"),
        ("🖥️", "DASHBOARD", "Peta interaktif + KPI + tabel per kecamatan"),
        ("✅", "VERIFIKASI", "Workflow staff enforcement + audit trail"),
    ]
    for i, (ico, title, desc) in enumerate(steps):
        x = Inches(0.5 + i * 2.5); y = Inches(2.2)
        card = add_bar(s, x, y, Inches(2.3), Inches(3.5), LIGHT_RGB)
        card.line.color.rgb = NAVY_RGB; card.line.width = Pt(1)
        # Icon
        add_text(s, ico, x, y + Inches(0.3), Inches(2.3), Inches(0.8), size=32, align=PP_ALIGN.CENTER)
        # Number badge
        badge = add_bar(s, x + Inches(0.8), y + Inches(1.2), Inches(0.7), Inches(0.4), NAVY_RGB)
        add_text(s, f"0{i+1}", x + Inches(0.8), y + Inches(1.2), Inches(0.7), Inches(0.4), size=14, bold=True, color=WHITE_RGB, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Title
        add_text(s, title, x + Inches(0.1), y + Inches(1.75), Inches(2.1), Inches(0.4), size=14, bold=True, color=NAVY_RGB, align=PP_ALIGN.CENTER)
        # Desc
        add_text(s, desc, x + Inches(0.15), y + Inches(2.2), Inches(2.0), Inches(1.2), size=10, color=GRAY_RGB, align=PP_ALIGN.CENTER)

    # ============ SLIDE 4: FITUR UTAMA ============
    s = prs.slides.add_slide(blank_layout); add_bg(s)
    add_header(s, "Fitur Utama", "03")
    fitur_cards = [
        ("Data Pipeline", "• 1,09M bangunan diproses\n• 31 polygon BPS\n• Spatial matching engine\n• Data quality audit"),
        ("Backend API", "• 7 REST endpoint\n• Multi-filter query\n• 2-tier caching\n• Cold 15s → warm 0,4s"),
        ("Peta Interaktif", "• Leaflet + satelit\n• 31 polygon kecamatan\n• Marker clustering\n• Filter 6-kategori"),
        ("Dashboard", "• 4 KPI cards realtime\n• Tabel per kecamatan\n• Stacked bar 3-segment\n• Drill-down interaktif"),
        ("Verifikasi", "• Panel 4-status\n• Bulk update API\n• Audit trail\n• Role-based access"),
        ("Production", "• Snapshot caching\n• Artisan refresh\n• Dokumentasi lengkap\n• SQL audit proof"),
    ]
    for i, (title, desc) in enumerate(fitur_cards):
        col = i % 3; row = i // 3
        x = Inches(0.4 + col * 4.3); y = Inches(1.3 + row * 2.8)
        card = add_bar(s, x, y, Inches(4.1), Inches(2.6), LIGHT_RGB)
        card.line.color.rgb = NAVY_RGB; card.line.width = Pt(1)
        # Title band
        band = add_bar(s, x, y, Inches(4.1), Inches(0.55), NAVY_RGB)
        add_text(s, title, x, y, Inches(4.1), Inches(0.55), size=16, bold=True, color=WHITE_RGB, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Content
        add_multiline(s, desc.split("\n"), x + Inches(0.2), y + Inches(0.7), Inches(3.8), Inches(1.8), size=12, color=NAVY_RGB)

    # ============ SLIDE 5: OUTPUT & OUTCOME ============
    s = prs.slides.add_slide(blank_layout); add_bg(s)
    add_header(s, "Output & Outcome", "04")

    # OUTPUT side
    card1 = add_bar(s, Inches(0.5), Inches(1.3), Inches(6), Inches(5.5), LIGHT_RGB)
    card1.line.color.rgb = NAVY_RGB; card1.line.width = Pt(1)
    add_bar(s, Inches(0.5), Inches(1.3), Inches(6), Inches(0.6), NAVY_RGB)
    add_text(s, "OUTPUT (Deliverable)", Inches(0.5), Inches(1.3), Inches(6), Inches(0.6), size=16, bold=True, color=WHITE_RGB, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    output_items = [
        "✓ Aplikasi web terdeploy di production DPUTR",
        "✓ Peta interaktif 31 polygon kecamatan",
        "✓ 97.000+ titik deteksi bangunan",
        "✓ Dashboard KPI + tabel breakdown",
        "✓ Workflow verifikasi + audit trail",
        "✓ Source code + dokumentasi lengkap",
        "✓ Training staff 1 batch",
    ]
    add_multiline(s, output_items, Inches(0.8), Inches(2.1), Inches(5.5), Inches(4.5), size=13, color=NAVY_RGB)

    # OUTCOME side
    card2 = add_bar(s, Inches(6.8), Inches(1.3), Inches(6), Inches(5.5), LIGHT_RGB)
    card2.line.color.rgb = SUCCESS_RGB; card2.line.width = Pt(1)
    add_bar(s, Inches(6.8), Inches(1.3), Inches(6), Inches(0.6), SUCCESS_RGB)
    add_text(s, "OUTCOME (Manfaat Terukur)", Inches(6.8), Inches(1.3), Inches(6), Inches(0.6), size=16, bold=True, color=WHITE_RGB, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    outcome_items = [
        "▲ Visibilitas 96.322 bangunan wajib PBG",
        "▲ Enforcement target langsung dipetakan",
        "▲ Basis data terpadu dengan PBB & tata ruang",
        "▲ Data-driven PAD policy untuk Bupati/DPRD",
        "▲ Dokumentasi kinerja DPUTR",
        "▲ ROI minimum 685× pada skenario realistis",
        "▲ Break-even cukup 15 PBG ditertibkan",
    ]
    add_multiline(s, outcome_items, Inches(7.1), Inches(2.1), Inches(5.5), Inches(4.5), size=13, color=NAVY_RGB)

    # ============ SLIDE 6: POTENSI PENDAPATAN ============
    s = prs.slides.add_slide(blank_layout); add_bg(s)
    add_header(s, "Potensi Pendapatan (ROI)", "05")
    add_text(s, "Rata-rata retribusi per PBG di Kab. Bandung = Rp 10,66 juta (data real SIMBG)", Inches(0.5), Inches(1.1), Inches(12), Inches(0.4), size=13, color=GRAY_RGB, align=PP_ALIGN.CENTER)

    # Scenario cards
    scenarios = [
        ("1%", "963", "Rp 10,3 M", GRAY_RGB),
        ("5%", "4.816", "Rp 51,3 M", GRAY_RGB),
        ("10%", "9.632", "Rp 102,7 M", SUCCESS_RGB),
        ("20%", "19.264", "Rp 205,4 M", SUCCESS_RGB),
        ("30%", "28.897", "Rp 308,0 M", GOLD_RGB),
    ]
    for i, (pct, pbg, pad, color) in enumerate(scenarios):
        x = Inches(0.5 + i * 2.53); y = Inches(1.8)
        card = add_bar(s, x, y, Inches(2.4), Inches(2.8), LIGHT_RGB)
        card.line.color.rgb = color; card.line.width = Pt(2)
        # Percentage badge top
        add_bar(s, x, y, Inches(2.4), Inches(0.7), color)
        add_text(s, pct + " enforcement", x, y, Inches(2.4), Inches(0.7), size=14, bold=True, color=WHITE_RGB, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, pbg + " PBG", x, y + Inches(0.9), Inches(2.4), Inches(0.5), size=14, color=GRAY_RGB, align=PP_ALIGN.CENTER)
        add_text(s, pad, x, y + Inches(1.5), Inches(2.4), Inches(0.8), size=22, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(s, "potensi PAD", x, y + Inches(2.3), Inches(2.4), Inches(0.4), size=10, color=GRAY_RGB, align=PP_ALIGN.CENTER)

    # Bottom: ROI highlight
    roi_box = add_bar(s, Inches(1), Inches(5.0), Inches(11.3), Inches(1.8), SUCCESS_RGB)
    add_text(s, f"INVESTASI {int(PRICE_TOTAL/1_000_000)} jt  →  ROI >685x  →  BREAK-EVEN cukup 15 PBG",
             Inches(1), Inches(5.0), Inches(11.3), Inches(1.8), size=22, bold=True, color=WHITE_RGB, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ============ SLIDE 7: BREAKDOWN HARGA ============
    s = prs.slides.add_slide(blank_layout); add_bg(s)
    add_header(s, "Investasi Anda", "06")

    pkg_boxes = [
        ("A", "Transfer, Deploy & Rekonsiliasi", f"Rp {PRICE_A_TRANSFER/1_000_000:.0f} jt", NAVY_RGB),
        ("B", "Pembangunan Fitur Satelit", f"Rp {PRICE_B_SATELIT/1_000_000:.0f} jt", NAVY_RGB),
        ("C", "Data Enrichment Geocoding (opsional)", f"Rp {PRICE_C_TRAINING/1_000_000:.0f} jt", GRAY_RGB),
        ("D", "Maintenance 6 Bulan (opsional)", f"Rp {PRICE_D_MAINT/1_000_000:.0f} jt", GRAY_RGB),
    ]
    for i, (code, desc, price, color) in enumerate(pkg_boxes):
        y = Inches(1.3 + i * 0.9)
        # Code badge
        add_bar(s, Inches(0.5), y, Inches(0.9), Inches(0.7), color)
        add_text(s, code, Inches(0.5), y, Inches(0.9), Inches(0.7), size=22, bold=True, color=WHITE_RGB, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Desc
        add_text(s, desc, Inches(1.6), y + Inches(0.1), Inches(8), Inches(0.5), size=16, color=NAVY_RGB, anchor=MSO_ANCHOR.MIDDLE)
        # Price
        add_text(s, price, Inches(9.8), y + Inches(0.1), Inches(3), Inches(0.5), size=18, bold=True, color=color, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    # Total
    y = Inches(5.2)
    add_bar(s, Inches(0.5), y, Inches(12.3), Inches(1.2), GOLD_RGB)
    add_text(s, "TOTAL INVESTASI (A+B+C+D)", Inches(0.5), y, Inches(6), Inches(1.2), size=20, bold=True, color=WHITE_RGB, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, f"Rp {PRICE_TOTAL/1_000_000:.0f} juta", Inches(6.5), y, Inches(6), Inches(1.2), size=32, bold=True, color=WHITE_RGB, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Disclaimer
    add_text(s, "Bandingkan dengan potensi PAD realistis: Rp 102,7 miliar (1000x return)",
             Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4), size=13, color=SUCCESS_RGB, align=PP_ALIGN.CENTER, bold=True)

    # ============ SLIDE 8: TIMELINE ============
    s = prs.slides.add_slide(blank_layout); add_bg(s)
    add_header(s, "Timeline 2,5 Bulan", "07")

    phases = [
        ("Transfer & Deploy", "M 1-2", "Production live", NAVY_RGB),
        ("Data Pipeline", "M 3-4", "1,09M row classified", NAVY_RGB),
        ("Backend API", "M 5-6", "7 endpoint aktif", NAVY_RGB),
        ("Frontend + Dashboard", "M 7-9", "Peta + filter live", NAVY_RGB),
        ("Verify & Testing", "M 10", "SQL audit proof", SUCCESS_RGB),
        ("UAT + Handover", "M 11", "Training + sign-off", GOLD_RGB),
    ]
    # Horizontal gantt-style
    step_w = Inches(1.95)
    for i, (title, week, deliv, color) in enumerate(phases):
        x = Inches(0.5 + i * 2.1); y = Inches(2.5)
        # Arrow shape
        shp = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW if i < len(phases)-1 else MSO_SHAPE.RECTANGLE, x, y, step_w, Inches(1.0))
        shp.fill.solid(); shp.fill.fore_color.rgb = color; shp.line.fill.background()
        add_text(s, week, x, y + Inches(0.15), step_w, Inches(0.4), size=13, bold=True, color=WHITE_RGB, align=PP_ALIGN.CENTER)
        add_text(s, title, x, y + Inches(0.55), step_w, Inches(0.4), size=10, color=WHITE_RGB, align=PP_ALIGN.CENTER)
        # Deliverable below
        add_text(s, deliv, x, y + Inches(1.2), step_w, Inches(0.5), size=10, color=GRAY_RGB, align=PP_ALIGN.CENTER, bold=True)

    # Bottom summary
    add_bar(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.2), LIGHT_RGB)
    add_text(s, "11 minggu · 4 payment milestone · Garansi 60 hari", Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.2), size=18, bold=True, color=NAVY_RGB, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ============ SLIDE 9: CALL TO ACTION ============
    s = prs.slides.add_slide(blank_layout); add_bg(s, NAVY_RGB)
    add_bar(s, 0, Inches(0.5), prs.slide_width, Inches(0.1), GOLD_RGB)

    add_text(s, "Siap Mulai?", Inches(0.5), Inches(1.2), Inches(12.3), Inches(1), size=48, bold=True, color=WHITE_RGB, align=PP_ALIGN.CENTER)
    add_text(s, f"Investasi Rp {PRICE_TOTAL/1_000_000:.0f} juta dengan potensi ROI minimum Rp 102 miliar", Inches(0.5), Inches(2.5), Inches(12.3), Inches(0.6), size=18, color=GOLD_RGB, align=PP_ALIGN.CENTER)

    # Key stats boxes
    stats_x = [Inches(1.0), Inches(5.1), Inches(9.2)]
    stats = [
        ("96.322", "Bangunan tanpa izin terdeteksi"),
        ("685×", "Minimum Return on Investment"),
        ("11 minggu", "Dari kickoff ke production"),
    ]
    for (x, (big, small)) in zip(stats_x, stats):
        y = Inches(3.5)
        box = add_bar(s, x, y, Inches(3.3), Inches(2.2), WHITE_RGB)
        box.line.color.rgb = GOLD_RGB; box.line.width = Pt(2)
        add_text(s, big, x, y + Inches(0.3), Inches(3.3), Inches(1), size=32, bold=True, color=NAVY_RGB, align=PP_ALIGN.CENTER)
        add_text(s, small, x, y + Inches(1.3), Inches(3.3), Inches(0.8), size=13, color=GRAY_RGB, align=PP_ALIGN.CENTER)

    add_text(s, f"Hubungi {VENDOR_NAME}", Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.5), size=16, color=WHITE_RGB, align=PP_ALIGN.CENTER, bold=True)
    add_text(s, "📧 email · 📱 WhatsApp · 🌐 portfolio", Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4), size=12, color=LIGHT_RGB, align=PP_ALIGN.CENTER)

    out = os.path.join(OUT_DIR, 'Presentasi_Monitoring_Satelit.pptx')
    prs.save(out)
    print(f"✓ PPTX: {out}")
    return out


if __name__ == "__main__":
    docx_path = build_docx()
    pdf_path = build_pdf()
    pptx_path = build_pptx()
    print("\nAll proposal docs generated in:", OUT_DIR)
