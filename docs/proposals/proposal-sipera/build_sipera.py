"""
Build profesional proposal docs (.docx, .pdf, .pptx) untuk SIPERA —
Sistem Perumahan Rakyat: Geospatial Intelligence Platform.

Mirror format proposal-satelit (build_proposal.py) supaya konsisten.
"""
import os
from datetime import datetime

# ====================== CONSTANTS ======================
OUT_DIR = os.path.dirname(os.path.abspath(__file__))
CLIENT_NAME = "DPUTR Kab. Bandung"
PROJECT_NAME = "SIPERA — Sistem Perumahan Rakyat"
PROJECT_SUBTITLE = "Geospatial Intelligence Platform untuk Pemetaan & Monitoring Perumahan Skala Kabupaten"
VENDOR_NAME = "Nauval Zulfikar, Fauzan Iraldi, Deri"
PJ_NAME = "Fauzan Iraldi"
PJ_ROLE = "Penanggung Jawab"
EXEC_NAME = "Nauval Zulfikar"
EXEC_ROLE = "Pelaksana"
DATE_STR = datetime.now().strftime("%d %B %Y")
PROPOSAL_NO = f"PROP-SIPERA-{datetime.now().strftime('%Y%m')}-001"

# ===== Variant configuration via env vars =====
# SIPERA_FEE_MODE: 'one_time' (default) | 'annual'
# SIPERA_MULTIPLIER: float, default 1.0 (e.g. 2.0 for 2X pricing)
FEE_MODE = os.environ.get('SIPERA_FEE_MODE', 'one_time')
PRICE_MULTIPLIER = float(os.environ.get('SIPERA_MULTIPLIER', '1.0'))

_BASE = {
    'TOTAL': 110_000_000,
    'A': 25_000_000,
    'B': 75_000_000,
    'C': 6_000_000,
    'D': 4_000_000,
}
PRICE_TOTAL = int(_BASE['TOTAL'] * PRICE_MULTIPLIER)
PRICE_A = int(_BASE['A'] * PRICE_MULTIPLIER)
PRICE_B = int(_BASE['B'] * PRICE_MULTIPLIER)
PRICE_C = int(_BASE['C'] * PRICE_MULTIPLIER)
PRICE_D = int(_BASE['D'] * PRICE_MULTIPLIER)

# Labels & filename suffix based on variant
PRICE_SUFFIX = " / tahun" if FEE_MODE == 'annual' else ""
FEE_LABEL = "Annual Subscription" if FEE_MODE == 'annual' else "One-Time"
FILE_TAG = []
if FEE_MODE == 'annual':
    FILE_TAG.append('Annual')
if PRICE_MULTIPLIER != 1.0:
    FILE_TAG.append(f'{PRICE_MULTIPLIER:g}X')
_TAG = ('_' + '_'.join(FILE_TAG)) if FILE_TAG else ''
OUTPUT_FILENAME_DOCX = f'Proposal_SIPERA{_TAG}.docx'
OUTPUT_FILENAME_PDF = f'Proposal_SIPERA{_TAG}.pdf'
OUTPUT_FILENAME_PPTX = f'Presentasi_SIPERA{_TAG}.pptx'

# Colors (brand) — same as satelit
NAVY = (28, 45, 84)
GOLD = (201, 162, 39)
GRAY = (90, 90, 90)
LIGHT = (245, 247, 250)
DANGER = (220, 53, 69)
SUCCESS = (40, 167, 69)


# ====================== DOCX ======================
def build_docx():
    from docx import Document
    from docx.shared import Pt, RGBColor, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_ALIGN_VERTICAL
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    doc = Document()
    for s in doc.sections:
        s.top_margin = Cm(2.2); s.bottom_margin = Cm(2.2)
        s.left_margin = Cm(2.5); s.right_margin = Cm(2.0)

    styles = doc.styles
    n = styles['Normal']
    n.font.name = 'Calibri'; n.font.size = Pt(11)

    def set_cell_bg(cell, rgb_hex):
        tc = cell._tc; tcPr = tc.get_or_add_tcPr()
        shd = OxmlElement('w:shd')
        shd.set(qn('w:val'), 'clear'); shd.set(qn('w:color'), 'auto'); shd.set(qn('w:fill'), rgb_hex)
        tcPr.append(shd)

    def H(text, level=1, color=NAVY, size=None, align=None, bold=True):
        sizes = {1: 22, 2: 16, 3: 13}
        p = doc.add_paragraph()
        if align == 'center': p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r = p.add_run(text)
        r.font.size = Pt(size or sizes.get(level, 13))
        r.font.bold = bold
        r.font.color.rgb = RGBColor(*color)
        return p

    def P(text, bold=False, italic=False, color=None, size=11, align=None, space_after=6):
        p = doc.add_paragraph()
        if align == 'center': p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_after = Pt(space_after)
        r = p.add_run(text)
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        if color: r.font.color.rgb = RGBColor(*color)
        return p

    def rupiah(num, with_suffix=False):
        base = f"Rp {num:,.0f}".replace(",", ".")
        return base + PRICE_SUFFIX if with_suffix else base

    # ============ COVER ============
    for _ in range(3): doc.add_paragraph()
    H("PROPOSAL PENAWARAN", level=1, size=36, align='center', color=NAVY)
    H(PROJECT_NAME, level=2, size=20, align='center', color=GOLD)
    P(PROJECT_SUBTITLE, align='center', italic=True, color=GRAY, size=12)
    doc.add_paragraph()
    P("Diajukan kepada:", align='center', color=GRAY, size=11)
    P(CLIENT_NAME, align='center', bold=True, size=14, color=NAVY)
    doc.add_paragraph(); doc.add_paragraph()

    tbl = doc.add_table(rows=3, cols=2); tbl.autofit = False
    tbl.columns[0].width = Cm(5); tbl.columns[1].width = Cm(10)
    for i, (lbl, val) in enumerate([("Nomor Proposal", PROPOSAL_NO),
                                     ("Tanggal", DATE_STR),
                                     ("Penyedia", VENDOR_NAME)]):
        tbl.cell(i, 0).text = lbl; tbl.cell(i, 1).text = val
        set_cell_bg(tbl.cell(i, 0), 'E8ECEF')

    doc.add_paragraph()
    P(f"Model Pembiayaan: {FEE_LABEL}", align='center', bold=True, color=GOLD, size=13)
    doc.add_paragraph()
    P('"Foundation geospasial production-grade: dari prototype 1-kabupaten', align='center', italic=True, color=NAVY, size=12)
    P('jadi platform scalable yang siap deploy ke 27 kabupaten Jawa Barat."', align='center', italic=True, color=NAVY, size=12)
    doc.add_page_break()

    # ============ EXECUTIVE SUMMARY ============
    H("RINGKASAN EKSEKUTIF", level=1)
    P("Sistem Monitoring Satelit yang sudah terdeploy di DPUTR Kab. Bandung berhasil memetakan 1,09 juta bangunan. "
      "Untuk naik level dari MVP ke platform production multi-kabupaten yang mampu menangani 4 juta+ bangunan "
      "dengan response sub-detik, dibutuhkan upgrade fondasi geospasial enterprise-grade. "
      "SIPERA adalah paket upgrade tersebut.")
    doc.add_paragraph()
    H("Angka Kunci", level=2)
    t = doc.add_table(rows=4, cols=2); t.style = 'Light Grid Accent 1'
    kv = [
        ("Total bangunan diproses (setelah upgrade)", "4,07 juta (3,7× dari MVP)"),
        ("Response time (cached)", "<300 ms (dari 22-30 detik)"),
        ("Kabupaten siap deploy", "Template ready (~1 hari per kab)"),
        ("Potensi PAD direct (10% realistis)", "Rp 383,8 miliar"),
    ]
    for i, (k, v) in enumerate(kv):
        t.cell(i, 0).text = k; t.cell(i, 1).text = v
        for c in t.rows[i].cells:
            c.vertical_alignment = WD_ALIGN_VERTICAL.CENTER
        set_cell_bg(t.cell(i, 0), 'E8ECEF')
    doc.add_paragraph()
    P("Dataset di-upgrade dari Microsoft Footprints (1,09M) ke kombinasi Microsoft + Google Open Buildings (4,07M), "
      "ditopang stack spatial database industri (PostGIS + Martin Vector Tile + Redis).", italic=True, color=GRAY)
    P("Artinya: dengan investasi Rp 110 juta, potensi tambahan PAD direct mencapai "
      "Rp 281 miliar — rasio ROI >2.800× pada skenario realistis.", bold=True, color=SUCCESS, size=12)
    doc.add_page_break()

    # ============ 1. LATAR BELAKANG ============
    H("1. LATAR BELAKANG", level=1)
    P("Sistem Monitoring Satelit Bangunan (proposal pendahulu, PROP-SATELIT-202604-001) "
      "telah berhasil membangun MVP enforcement bangunan tanpa izin di Kab. Bandung dengan 1,09 juta data "
      "Microsoft Footprints. Setelah implementasi, teridentifikasi tiga kebutuhan upgrade fondamental:")
    for b in [
        "Cakupan data Microsoft (1,09M) belum cover seluruh bangunan; Google Open Buildings menambah ~3M polygon yang sebelumnya tidak terdeteksi.",
        "Stack MariaDB + Leaflet geojson tidak scalable untuk 4M+ polygon — render bisa makan 22-30 detik (cold).",
        "Sistem belum bisa di-replikasi ke kabupaten lain tanpa rebuild dari nol — tidak ada template multi-tenant.",
    ]:
        doc.add_paragraph(b, style='List Bullet')
    doc.add_paragraph()
    P("SIPERA adalah paket upgrade yang membawa fondasi sistem ke level production enterprise: "
      "PostGIS spatial database, Martin Vector Tile Server (industry-standard MVT), Redis caching layer, "
      "server-side aggregation, dan filter-gated lazy rendering. Hasilnya: platform geospasial scalable "
      "yang siap untuk 27 kabupaten Jawa Barat.")
    doc.add_page_break()

    # ============ 2. OUTPUT & OUTCOME ============
    H("2. OUTPUT & OUTCOME", level=1)
    H("Output — Yang DPUTR Terima", level=2)
    for b in [
        "Spatial database PostGIS 16 (Docker, tuned) berisi 4,07 juta polygon bangunan ter-indeks GIST.",
        "Martin Vector Tile Server (MVT) untuk rendering peta level Mapbox/ESRI.",
        "Redis caching layer dengan 24h TTL + auto-prewarm cron harian.",
        "Server-side spatial aggregation endpoint (cluster z6-13) — render <300ms untuk 4M+ polygon.",
        "Filter-gated rendering UX (choropleth default → drill-down per kecamatan/desa).",
        "Lazy rendering: pan debounce, centroid-first, viewport clipping.",
        "13 ranked performance optimization (load test report + before/after benchmark).",
        "Template deploy multi-kabupaten (per-kab config, ~1 hari setup).",
        "Source code + dokumentasi spatial stack + runbook deploy/migrasi.",
    ]:
        doc.add_paragraph(b, style='List Bullet')

    H("Outcome — Manfaat Terukur", level=2)
    for b in [
        "Coverage data naik 3,7× (1,09M → 4,07M bangunan) — tidak ada blind spot.",
        "User experience sub-detik untuk peta 4M polygon (sebelumnya 22-30 detik cold).",
        "Foundation siap di-extend untuk fitur lanjutan: housing program, slum mapping, RDTR compliance, BPJS Perumahan, dll.",
        "Per-kabupaten deployment cost turun drastis — 1 hari setup vs 11 minggu rebuild dari nol.",
        "Open architecture, no vendor lock-in (PostGIS + MVT = standar industri).",
    ]:
        doc.add_paragraph(b, style='List Bullet')
    doc.add_page_break()

    # ============ 3. POTENSI PENDAPATAN (NEW) ============
    H("3. POTENSI PENDAPATAN (Kabupaten Bandung)", level=1)
    P("Section ini menghitung potensi pendapatan asli daerah (PAD) untuk DPUTR Kab. Bandung "
      "setelah implementasi SIPERA, dengan semua asumsi dijelaskan sumbernya supaya bisa "
      "diaudit ulang oleh tim DPUTR atau auditor independen.", italic=True)
    doc.add_paragraph()

    H("3.1 Sumber Data & Referensi", level=2)
    t = doc.add_table(rows=1, cols=2); t.style = 'Light Grid Accent 1'
    hdr = t.rows[0].cells
    hdr[0].text = "Data Point"; hdr[1].text = "Sumber"
    for c in hdr:
        set_cell_bg(c, '1C2D54')
        for p in c.paragraphs:
            for r in p.runs: r.font.color.rgb = RGBColor(255, 255, 255); r.font.bold = True
    for rd in [
        ("Penduduk & wilayah Kab. Bandung", "BPS Kab. Bandung 2023-2024"),
        ("PBG terbit Kab. Bandung", "SIMBG Kementerian PUPR (5.241 record real)"),
        ("Bangunan terdeteksi 1,09M", "Microsoft GlobalMLBuildingFootprints v3 (2024)"),
        ("Bangunan terdeteksi 4,07M", "Google Open Buildings v3 (2024) — release Jan 2024"),
        ("Aturan PBG & retribusi", "PP 16/2021 + Perda Kab. Bandung tentang Retribusi PBG"),
        ("Benchmark enforcement digital", "Studi tax compliance Direktorat Jenderal Pajak"),
    ]:
        row = t.add_row().cells
        for i, v in enumerate(rd): row[i].text = v
    doc.add_paragraph()

    H("3.2 Asumsi Kunci & Justifikasi Angka", level=2)
    t = doc.add_table(rows=1, cols=3); t.style = 'Light Grid Accent 1'
    hdr = t.rows[0].cells
    hdr[0].text = "Asumsi"; hdr[1].text = "Nilai"; hdr[2].text = "Sumber / Dasar Perhitungan"
    for c in hdr:
        set_cell_bg(c, '1C2D54')
        for p in c.paragraphs:
            for r in p.runs: r.font.color.rgb = RGBColor(255, 255, 255); r.font.bold = True
    for rd in [
        ("Penduduk Kab. Bandung", "3,82 juta jiwa", "BPS Kab. Bandung 2023"),
        ("Luas wilayah", "1.762 km²", "BPS Kab. Bandung"),
        ("Total kecamatan", "31 kecamatan", "BPS — 270 desa + 10 kelurahan"),
        ("Bangunan terdeteksi MVP", "1,09 juta", "Microsoft Footprints v3 (2024)"),
        ("Bangunan terdeteksi SIPERA", "4,07 juta", "Microsoft + Google Open Buildings v3"),
        ("Filter ≥200 m² tanpa izin (MVP)", "96.322 unit (8,8%)", "Hasil filter audit data project"),
        ("Estimasi target SIPERA", "~360.000 unit", "Asumsi proporsi 8,8% × 4,07M"),
        ("Retribusi rata-rata PBG", "Rp 10,66 juta", "Real data SIMBG 5.241 PBG = Rp 55,86 M total"),
        ("Enforcement realistis 10%", "36.000 PBG", "Best practice tax digital compliance"),
    ]:
        row = t.add_row().cells
        for i, v in enumerate(rd): row[i].text = v
    doc.add_paragraph()
    P("Validasi cross-check: 3,82 juta penduduk ÷ rata-rata 4 jiwa per rumah = 955.000 unit "
      "rumah tangga (range plausible terhadap 4,07M bangunan total termasuk non-residensial).",
      italic=True, color=GRAY, size=10)
    doc.add_paragraph()

    H("3.3 Skenario Potensi PAD dari 360.000 Bangunan Target", level=2)
    t = doc.add_table(rows=1, cols=4); t.style = 'Light Grid Accent 1'
    hdr = t.rows[0].cells
    hdr[0].text = "Skenario"; hdr[1].text = "Tingkat Enforcement"
    hdr[2].text = "PBG Tertarik"; hdr[3].text = "Potensi PAD"
    for c in hdr:
        set_cell_bg(c, '1C2D54')
        for p in c.paragraphs:
            for r in p.runs: r.font.color.rgb = RGBColor(255, 255, 255); r.font.bold = True
    for rd in [
        ("Super konservatif", "1%", "3.600", "Rp 38,4 miliar"),
        ("Konservatif", "5%", "18.000", "Rp 191,9 miliar"),
        ("Realistis", "10%", "36.000", "Rp 383,8 miliar"),
        ("Target wajar", "20%", "72.000", "Rp 767,5 miliar"),
        ("Agresif", "30%", "108.000", "Rp 1.151 miliar"),
    ]:
        row = t.add_row().cells
        for i, v in enumerate(rd): row[i].text = v
    doc.add_paragraph()
    P("Formula: PBG tertarik = 360.000 × tingkat enforcement. Potensi PAD = PBG tertarik × "
      "Rp 10,66 juta. Skenario realistis (10%) sebanding dengan baseline tax compliance "
      "improvement studies DJP setelah digitalisasi (8-15%).", italic=True, color=GRAY, size=10)
    doc.add_paragraph()
    P("Tambahan dari MVP Satelit (10% realistis): Rp 383,8 M vs Rp 102,7 M = "
      "+Rp 281 miliar PAD tambahan yang baru terjangkau setelah upgrade SIPERA.",
      bold=True, color=SUCCESS, size=11)
    doc.add_paragraph()

    H("3.4 Justifikasi Tingkat Enforcement", level=2)
    t = doc.add_table(rows=1, cols=3); t.style = 'Light Grid Accent 1'
    hdr = t.rows[0].cells
    hdr[0].text = "Tingkat"; hdr[1].text = "Kondisi"; hdr[2].text = "Basis Argumentasi"
    for c in hdr:
        set_cell_bg(c, '1C2D54')
        for p in c.paragraphs:
            for r in p.runs: r.font.color.rgb = RGBColor(255, 255, 255); r.font.bold = True
    for rd in [
        ("1%", "Worst-case adopsi minim", "Tanpa kampanye, no follow-up"),
        ("5%", "Adopsi natural", "Tingkat compliance suka rela properti rata-rata"),
        ("10%", "Workflow enforcement jalan", "DJP benchmark digital tax improvement 8-15%"),
        ("20%", "Dedicated team enforcement", "Sebanding capaian Kota Bekasi/Tangerang"),
        ("30%", "Kampanye + sanksi penuh", "Best practice pemkot dengan SK Bupati"),
    ]:
        row = t.add_row().cells
        for i, v in enumerate(rd): row[i].text = v
    doc.add_paragraph()

    H("3.5 Productivity Gain Staff Enforcement", level=2)
    t = doc.add_table(rows=1, cols=3); t.style = 'Light Grid Accent 1'
    hdr = t.rows[0].cells
    hdr[0].text = "Aktivitas"; hdr[1].text = "Sebelum"; hdr[2].text = "Setelah SIPERA"
    for c in hdr:
        set_cell_bg(c, '1C2D54')
        for p in c.paragraphs:
            for r in p.runs: r.font.color.rgb = RGBColor(255, 255, 255); r.font.bold = True
    for rd in [
        ("Audit 1 kecamatan", "15-20 menit", "5-8 menit (3× cepat)"),
        ("Buka dashboard", "22-30 detik", "<300 ms"),
        ("Pan/zoom peta", "Janky 15 fps", "Smooth 50 fps"),
        ("User concurrent", "5-10 staff", "50-80 staff"),
        ("Akses HP lapangan", "Crash", "Lancar"),
    ]:
        row = t.add_row().cells
        for i, v in enumerate(rd): row[i].text = v
    doc.add_paragraph()
    P("Throughput audit 3× lipat → realisasi PAD lebih cepat. Asumsi 31 kecamatan × 8 audit "
      "per bulan × 12 bulan × (15 menit hemat) = ~750 jam staff/tahun untuk seluruh DPUTR.",
      italic=True, color=GRAY, size=10)
    doc.add_paragraph()

    H("3.6 Ringkasan ROI SIPERA", level=2)
    t_roi = doc.add_table(rows=5, cols=2); t_roi.style = 'Light Grid Accent 1'
    roi_data = [
        ("Investasi SIPERA (Paket A+B)", "Rp 100 juta"),
        ("Tambahan PAD direct (skenario realistis 10%)", "+Rp 281 miliar"),
        ("Rasio ROI direct", ">2.800×"),
        ("Break-even", "10 PBG baru ditertibkan"),
        ("Total ROI gabungan (Satelit+SIPERA Rp 260jt)", "~1.476× terhadap PAD Rp 383,8 M"),
    ]
    for i, (k, v) in enumerate(roi_data):
        t_roi.cell(i, 0).text = k; t_roi.cell(i, 1).text = v
        set_cell_bg(t_roi.cell(i, 0), 'E8ECEF')
    doc.add_paragraph()
    P("Semua angka di atas adalah tambahan dari proposal Satelit Rp 150 juta. "
      "Total ROI dihitung terhadap PAD skenario realistis Kab. Bandung saja "
      "(belum termasuk denda administratif, PBB, atau efek tata ruang).",
      italic=True, color=NAVY, size=10)
    doc.add_page_break()

    # ============ 4. RUANG LINGKUP ============
    H("4. RUANG LINGKUP PEKERJAAN", level=1)
    P("Seluruh scope pekerjaan disajikan dalam bentuk tabel untuk kemudahan audit dan tracking deliverable. "
      "Rincian nilai investasi ada pada Bagian 5.", italic=True, color=GRAY)
    doc.add_paragraph()

    def pkg_table_nopr(code, name, items):
        rows = [["Kode Paket", code], ["Nama Paket", name]]
        t = doc.add_table(rows=len(rows) + 1 + len(items), cols=2)
        t.style = 'Light Grid Accent 1'
        t.columns[0].width = Cm(5); t.columns[1].width = Cm(11)
        for i, (k, v) in enumerate(rows):
            t.cell(i, 0).text = k; t.cell(i, 1).text = v
            set_cell_bg(t.cell(i, 0), '1C2D54')
            for p in t.cell(i, 0).paragraphs:
                for r in p.runs:
                    r.font.color.rgb = RGBColor(255, 255, 255); r.font.bold = True
        hdr_idx = len(rows)
        merged = t.cell(hdr_idx, 0).merge(t.cell(hdr_idx, 1))
        merged.text = "Deliverable & Fitur Tercakup"
        set_cell_bg(merged, 'C9A227')
        for p in merged.paragraphs:
            for r in p.runs:
                r.font.color.rgb = RGBColor(255, 255, 255); r.font.bold = True
        for j, item in enumerate(items):
            row_idx = hdr_idx + 1 + j
            t.cell(row_idx, 0).text = f"{code}.{j+1}"
            t.cell(row_idx, 1).text = item
            set_cell_bg(t.cell(row_idx, 0), 'E8ECEF')
            for p in t.cell(row_idx, 0).paragraphs:
                for r in p.runs:
                    r.font.bold = True; r.font.color.rgb = RGBColor(*NAVY)
        doc.add_paragraph()

    # ===== PAKET A =====
    H("Paket A — Foundation Upgrade: Migration & Spatial Stack Setup", level=2)
    pkg_table_nopr("A", "Migrasi MariaDB → PostGIS + Setup Spatial Infrastructure", [
        "Setup PostGIS 16 (Docker) dengan tuning shared_buffers=2GB, work_mem=64MB, maintenance_work_mem=512MB.",
        "GIST spatial indexing untuk 4M+ geometries (rendering p95 <500ms).",
        "Migrasi data MariaDB → PostGIS dengan MultiPolygon + Polygon dual-handling (ST_GeometryN).",
        "Setup Redis caching layer (24h TTL + auto-invalidation pada verify event).",
        "Daily prewarm cron (Windows Task Scheduler / systemd) untuk warm cache.",
        "Reproducible deploy script (Docker Compose) + rollback runbook.",
    ])

    # ===== PAKET B (big, 7 sub-paket) =====
    H("Paket B — Geospatial Engineering Platform", level=2)
    P("Paket B dipecah menjadi 7 sub-paket teknis:", italic=True)

    subpkg = [
        ("B.1", "Multi-Source Data Pipeline (4M+ rows)", [
            "Ingest Google Open Buildings 4,07 juta polygon (3,7× dari Microsoft).",
            "Multi-source merger: Microsoft + Google + manual annotation.",
            "Kecamatan enrichment 4,07M rows via Python + scipy KDTree.",
            "High-performance spatial matching engine (~36K rows/sec).",
            "Data quality audit + reconciliation report multi-source.",
            "Reproducible pipeline scripts (Python + Bash + Artisan).",
        ]),
        ("B.2", "Enterprise Spatial Database (PostGIS + Redis)", [
            "PostGIS 16 schema dengan multi-source tables + indexed geometry.",
            "Custom MVT SQL function (zoom-aware, modulo sampling untuk low-zoom).",
            "Redis cache layer dengan 24h TTL untuk tile + cluster endpoints.",
            "MultiPolygon ST_GeometryN handling untuk irregular shapes.",
            "PostgreSQL tuning + benchmark report (cold vs warm).",
        ]),
        ("B.3", "Vector Tile Rendering Pipeline (Mapbox-grade)", [
            "Martin Vector Tile Server deployment (industry-standard MVT).",
            "TilesController endpoint dengan tile cache headers (max-age=86400).",
            "Leaflet.VectorGrid integration (replace GeoJSON layer untuk 4M polygon).",
            "5-state polygon palette aligned dengan cluster pins (terbit/proses/ditolak/tanpa-izin/PBG).",
            "Modulo sampling (b.id % 5) untuk low-zoom z≤15 supaya render tetap cepat.",
        ]),
        ("B.4", "Server-Side Spatial Aggregation", [
            "Endpoint /clusters dengan PostGIS grid bucketing (z6-13).",
            "Adaptive zoom-based aggregation (step = 112.5 / 2^zoom).",
            "24h cluster cache + daily prewarm via cron.",
            "Frontend mode-switch otomatis (cluster ↔ vector tile ↔ individual).",
        ]),
        ("B.5", "Filter-Gated Rendering UX", [
            "Choropleth default view: 31 polygon kecamatan dengan warna density.",
            "Materialized view kecamatan_violation_stats (daily refresh).",
            "Mode otomatis: no filter → choropleth, filter kecamatan → cluster+tile, filter desa → polygon individual.",
            "Empty-state UX yang informatif (bukan layar kosong).",
        ]),
        ("B.6", "Lazy Rendering Optimizations", [
            "Pan/zoom debounce 300ms (gak refetch tiap drag).",
            "Centroid-first rendering z14-15 (titik dulu, polygon z≥16).",
            "Progressive viewport clipping (drop polygon di luar viewport ketat).",
            "Loading skeleton + abort-on-pan race handling.",
        ]),
        ("B.7", "Performance Engineering & Production Readiness", [
            "13 ranked performance optimizations (load test + before/after report).",
            "Playwright-based benchmark dengan network timing breakdown.",
            "Tile feature LIMIT 10K + zoom-conditional sampling.",
            "Verified visual regression: warna polygon/cluster alignment.",
            "Health-check endpoint + monitoring runbook.",
        ]),
    ]
    total_rows = 1 + sum(1 + len(items) for _, _, items in subpkg)
    tb = doc.add_table(rows=total_rows, cols=3); tb.style = 'Light Grid Accent 1'
    tb.columns[0].width = Cm(2); tb.columns[1].width = Cm(5); tb.columns[2].width = Cm(9)
    hdr = tb.rows[0].cells
    hdr[0].text = "Kode"; hdr[1].text = "Sub-Paket"; hdr[2].text = "Deliverable"
    for c in hdr:
        set_cell_bg(c, '1C2D54')
        for p in c.paragraphs:
            for r in p.runs:
                r.font.color.rgb = RGBColor(255, 255, 255); r.font.bold = True
    ri = 1
    for code, name, items in subpkg:
        row = tb.rows[ri].cells
        row[0].text = code; row[1].text = name; row[2].text = ""
        for c in row:
            set_cell_bg(c, 'C9A227')
            for p in c.paragraphs:
                for r in p.runs:
                    r.font.bold = True; r.font.color.rgb = RGBColor(255, 255, 255)
        ri += 1
        for j, item in enumerate(items):
            row = tb.rows[ri].cells
            row[0].text = f"{code}.{j+1}"; row[1].text = ""; row[2].text = item
            set_cell_bg(row[0], 'E8ECEF')
            for p in row[0].paragraphs:
                for r in p.runs:
                    r.font.bold = True; r.font.color.rgb = RGBColor(*NAVY)
            ri += 1
    doc.add_paragraph()

    # ===== PAKET C =====
    H("Paket C — Multi-Kabupaten Templating (Opsional)", level=2)
    pkg_table_nopr("C", "Template Deploy ke Kabupaten Lain (1 hari setup)", [
        "Schema generic (tidak hardcode Kab. Bandung).",
        "Per-kabupaten config file (admin boundary, bbox, prewarm targets).",
        "Documentation deploy ke kabupaten baru (~1 hari setup).",
        "Tested deploy script untuk minimal 1 kabupaten referensi.",
    ])

    # ===== PAKET D =====
    H("Paket D — Maintenance & SLA 6 Bulan (Opsional)", level=2)
    pkg_table_nopr("D", "Maintenance & Support 6 Bulan", [
        "Bug-fix unlimited untuk defect spatial stack.",
        "Performance regression monitoring bulanan.",
        "Cache + prewarm health check.",
        "Response time 24 jam untuk issue blocker.",
    ])
    doc.add_page_break()

    # ============ 5. INVESTASI ============
    H("5. INVESTASI & PEMBAYARAN", level=1)
    P("Setelah melihat output, outcome, dan potensi PAD, berikut rincian investasi yang dibutuhkan:",
      italic=True, color=GRAY)
    doc.add_paragraph()

    # Highlight contrast box
    t_contrast = doc.add_table(rows=2, cols=2); t_contrast.style = 'Light Grid Accent 1'
    t_contrast.cell(0, 0).text = "Potensi PAD (realistis 10%)"
    t_contrast.cell(0, 1).text = "Rp 383,8 miliar"
    fee_short = "Annual" if FEE_MODE == 'annual' else "One-Time"
    t_contrast.cell(1, 0).text = f"Biaya SIPERA ({fee_short})"
    t_contrast.cell(1, 1).text = rupiah(PRICE_TOTAL, with_suffix=True)
    set_cell_bg(t_contrast.cell(0, 0), '28A745'); set_cell_bg(t_contrast.cell(0, 1), '28A745')
    set_cell_bg(t_contrast.cell(1, 0), '1C2D54'); set_cell_bg(t_contrast.cell(1, 1), '1C2D54')
    for r in t_contrast.rows:
        for c in r.cells:
            for p in c.paragraphs:
                for run in p.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)
                    run.font.bold = True
                    run.font.size = Pt(13)
    doc.add_paragraph()
    P("ROI direct >2.800× · Break-even: cukup 10 PBG baru ditertibkan dari 360.000 bangunan terdeteksi.",
      bold=True, color=SUCCESS, align='center', size=12)
    doc.add_paragraph()

    H("Rincian Harga per Paket", level=2)
    t = doc.add_table(rows=1, cols=3); t.style = 'Light Grid Accent 1'
    hdr = t.rows[0].cells
    hdr[0].text = "Paket"; hdr[1].text = "Deskripsi"; hdr[2].text = "Harga"
    for c in hdr:
        set_cell_bg(c, '1C2D54')
        for p in c.paragraphs:
            for r in p.runs:
                r.font.color.rgb = RGBColor(255, 255, 255); r.font.bold = True

    pkgs = [
        ("A", "Foundation Upgrade: Migration & Spatial Stack", rupiah(PRICE_A)),
        ("B", "Geospatial Engineering Platform (7 sub-paket)", rupiah(PRICE_B)),
        ("C", "Multi-Kabupaten Templating (opsional)", rupiah(PRICE_C)),
        ("D", "Maintenance & SLA 6 Bulan (opsional)", rupiah(PRICE_D)),
    ]
    for pkg in pkgs:
        row = t.add_row().cells
        for i, v in enumerate(pkg): row[i].text = v
    total_row = t.add_row().cells
    total_row[0].text = ""; total_row[1].text = f"TOTAL (A+B+C+D){' / tahun' if FEE_MODE == 'annual' else ''}"; total_row[2].text = rupiah(PRICE_TOTAL, with_suffix=True)
    for c in total_row:
        set_cell_bg(c, 'C9A227')
        for p in c.paragraphs:
            for r in p.runs: r.font.bold = True

    doc.add_paragraph()
    H("Term Pembayaran", level=2)
    t = doc.add_table(rows=1, cols=3); t.style = 'Light Grid Accent 1'
    hdr = t.rows[0].cells
    hdr[0].text = "Milestone"; hdr[1].text = "%"; hdr[2].text = "Nominal"
    for c in hdr:
        set_cell_bg(c, '1C2D54')
        for p in c.paragraphs:
            for r in p.runs:
                r.font.color.rgb = RGBColor(255, 255, 255); r.font.bold = True
    if FEE_MODE == 'annual':
        ms_list = [
            ("Pembayaran Tahun 1 (kickoff + onboarding)", "100%", rupiah(PRICE_TOTAL)),
            ("Renewal Tahun 2", "100%", rupiah(PRICE_TOTAL)),
            ("Renewal Tahun 3 (opsional)", "100%", rupiah(PRICE_TOTAL)),
        ]
    else:
        ms_list = [
            ("Down Payment (kickoff)", "30%", rupiah(int(PRICE_TOTAL * 0.3))),
            ("Paket A selesai (PostGIS+Redis live)", "20%", rupiah(int(PRICE_TOTAL * 0.2))),
            ("Paket B.1-B.4 selesai (MVT + cluster live)", "30%", rupiah(int(PRICE_TOTAL * 0.3))),
            ("Handover + UAT + benchmark report", "20%", rupiah(int(PRICE_TOTAL * 0.2))),
        ]
    for ms in ms_list:
        row = t.add_row().cells
        for i, v in enumerate(ms): row[i].text = v

    doc.add_paragraph()
    P("Catatan:", bold=True)
    base_notes_docx = [
        "Harga sudah termasuk source code + dokumentasi spatial stack + runbook deploy.",
        "Stand-alone proposal — bisa di-bundle dengan proposal Sistem Monitoring Satelit (Rp 150jt) dengan diskon sinergi 5-10jt.",
        "Deployment di infrastruktur klien (tidak include biaya server/domain).",
        "Penawaran berlaku 30 hari sejak tanggal diterbitkan.",
    ]
    if FEE_MODE == 'annual':
        notes_docx = base_notes_docx[:3] + [
            "Annual subscription mencakup: maintenance unlimited bug-fix, performance regression monitoring, security update, minor enhancement (jam kerja terbatas).",
            "Auto-renewal opsional tiap tahun; bisa di-cancel dengan notice 60 hari sebelum periode habis.",
            "SLA response 24 jam untuk issue blocker; uptime + cache health check inklusif.",
        ] + base_notes_docx[3:]
    else:
        notes_docx = base_notes_docx[:3] + ["Garansi bug-fix 60 hari setelah handover untuk spatial stack."] + base_notes_docx[3:]
    for b in notes_docx:
        doc.add_paragraph(b, style='List Bullet')
    doc.add_page_break()

    # ============ 6. TIMELINE ============
    H("6. TIMELINE", level=1)
    t = doc.add_table(rows=1, cols=3); t.style = 'Light Grid Accent 1'
    hdr = t.rows[0].cells
    hdr[0].text = "Fase"; hdr[1].text = "Minggu"; hdr[2].text = "Deliverable"
    for c in hdr:
        set_cell_bg(c, '1C2D54')
        for p in c.paragraphs:
            for r in p.runs:
                r.font.color.rgb = RGBColor(255, 255, 255); r.font.bold = True
    for ph in [
        ("Foundation Upgrade (PostGIS + Redis)", "1", "Spatial DB + cache layer live."),
        ("Data Pipeline 4M (Google + matching)", "2", "4,07M row terklasifikasi + indexed."),
        ("Vector Tile Server (Martin MVT)", "3", "MVT endpoint live + Leaflet.VectorGrid."),
        ("Server-Side Clustering", "4", "/clusters endpoint + frontend mode-switch."),
        ("Filter-Gated + Lazy Rendering", "5", "Choropleth + lazy UX."),
        ("Performance Engineering 13 ranks", "6", "Load test report + before/after benchmark."),
        ("Multi-Kab + UAT + Handover", "7", "Template deploy + dokumentasi + sign-off."),
    ]:
        row = t.add_row().cells
        for i, v in enumerate(ph): row[i].text = v

    doc.add_paragraph()
    P("Total durasi: 7 minggu (~1,75 bulan).", bold=True, color=NAVY)
    P("Catatan: Lebih singkat dari proposal satelit (11 minggu) karena infrastruktur dasar "
      "sudah tersedia dari proyek pendahulu.", italic=True, color=GRAY, size=10)
    doc.add_page_break()

    # ============ 7. PENUTUP ============
    H("7. PENUTUP", level=1)
    P("Proposal SIPERA menawarkan upgrade fondasi geospasial dari MVP enforcement (proposal satelit) "
      "ke platform production multi-kabupaten yang scalable. Dengan investasi Rp 110 juta, "
      "potensi tambahan PAD direct Kab. Bandung mencapai Rp 281 miliar (skenario realistis 10%) — "
      "ROI minimum 2.800× pada skenario realistis.")
    doc.add_paragraph()
    P("SIPERA bisa berdiri sendiri sebagai produk baru, atau di-bundle dengan proposal Sistem "
      "Monitoring Satelit sebagai paket lengkap. Output paket ini juga reusable untuk produk turunan: "
      "Sistem Monitoring Perumahan Subsidi, Slum Mapping, Tax Mapping, dan inisiatif geospasial DPUTR lainnya.")
    doc.add_paragraph(); doc.add_paragraph()
    P("Terima kasih atas kesempatan ini.", align='center', italic=True)
    doc.add_paragraph(); doc.add_paragraph()

    t = doc.add_table(rows=1, cols=2); t.autofit = False
    t.columns[0].width = Cm(8); t.columns[1].width = Cm(8)
    left_block = (
        f"{PJ_ROLE},\n\n\n\n"
        f"{PJ_NAME}\n\n\n"
        f"{EXEC_ROLE},\n\n\n\n"
        f"{EXEC_NAME}"
    )
    right_block = f"{CLIENT_NAME},\n\n\n\n__________________________"
    t.cell(0, 0).text = left_block
    t.cell(0, 1).text = right_block

    out = os.path.join(OUT_DIR, OUTPUT_FILENAME_DOCX)
    doc.save(out)
    print(f"✓ DOCX: {out}")
    return out


# ====================== PDF ======================
def build_pdf():
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm, mm
    from reportlab.lib.colors import HexColor, white
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, KeepTogether
    from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY

    NAVY_HEX = HexColor('#1C2D54'); GOLD_HEX = HexColor('#C9A227'); LIGHT_HEX = HexColor('#E8ECEF')
    GRAY_HEX = HexColor('#5A5A5A'); SUCCESS_HEX = HexColor('#28A745')

    out = os.path.join(OUT_DIR, OUTPUT_FILENAME_PDF)
    doc = SimpleDocTemplate(out, pagesize=A4, leftMargin=2.5 * cm, rightMargin=2 * cm,
                            topMargin=2.2 * cm, bottomMargin=2 * cm)
    styles = getSampleStyleSheet()
    h1 = ParagraphStyle('h1', parent=styles['Heading1'], fontSize=22, textColor=NAVY_HEX,
                        spaceAfter=14, leading=26, fontName='Helvetica-Bold')
    h2 = ParagraphStyle('h2', parent=styles['Heading2'], fontSize=15, textColor=NAVY_HEX,
                        spaceAfter=8, fontName='Helvetica-Bold')
    h3 = ParagraphStyle('h3', parent=styles['Heading3'], fontSize=12, textColor=NAVY_HEX,
                        spaceAfter=6, fontName='Helvetica-Bold')
    body = ParagraphStyle('body', parent=styles['Normal'], fontSize=11, leading=15,
                          spaceAfter=6, fontName='Helvetica', alignment=TA_JUSTIFY)
    bullet = ParagraphStyle('bullet', parent=body, leftIndent=18, bulletIndent=6, spaceAfter=3)
    center = ParagraphStyle('center', parent=body, alignment=TA_CENTER)
    cover_title = ParagraphStyle('ct', fontSize=32, alignment=TA_CENTER, textColor=NAVY_HEX,
                                  fontName='Helvetica-Bold', leading=40, spaceAfter=10)
    cover_sub = ParagraphStyle('cs', fontSize=20, alignment=TA_CENTER, textColor=GOLD_HEX,
                                fontName='Helvetica-Bold', leading=24, spaceAfter=8)
    cover_subtitle = ParagraphStyle('cst', fontSize=12, alignment=TA_CENTER, textColor=GRAY_HEX,
                                     fontName='Helvetica-Oblique', leading=16, spaceAfter=16)
    cover_tag = ParagraphStyle('ctag', fontSize=13, alignment=TA_CENTER, textColor=NAVY_HEX,
                                fontName='Helvetica-Oblique', leading=18)
    quote_style = ParagraphStyle('q', parent=body, alignment=TA_CENTER, fontSize=12,
                                  textColor=NAVY_HEX, fontName='Helvetica-Oblique')

    def rupiah(n, with_suffix=False):
        base = f"Rp {n:,.0f}".replace(",", ".")
        return base + PRICE_SUFFIX if with_suffix else base

    from reportlab.lib.enums import TA_LEFT
    cell_style = ParagraphStyle('cell', fontSize=10, leading=12.5, fontName='Helvetica',
                                 alignment=TA_LEFT)
    cell_style_b = ParagraphStyle('cell_b', fontSize=10, leading=12.5, fontName='Helvetica-Bold',
                                   alignment=TA_LEFT)

    def wp(text, bold=False):
        """Wrap cell text in Paragraph so it auto-wraps within column width."""
        return Paragraph(text, cell_style_b if bold else cell_style)

    story = []

    # COVER
    story.append(Spacer(1, 3 * cm))
    story.append(Paragraph("PROPOSAL PENAWARAN", cover_title))
    story.append(Paragraph(PROJECT_NAME, cover_sub))
    story.append(Paragraph(PROJECT_SUBTITLE, cover_subtitle))
    story.append(Spacer(1, 1 * cm))
    story.append(Paragraph("Diajukan kepada:", cover_tag))
    story.append(Spacer(1, 4 * mm))
    story.append(Paragraph(f"<b>{CLIENT_NAME}</b>",
                            ParagraphStyle('cn', fontSize=16, alignment=TA_CENTER,
                                           textColor=NAVY_HEX, fontName='Helvetica-Bold')))
    story.append(Spacer(1, 2 * cm))

    meta = [
        ["Nomor Proposal", PROPOSAL_NO],
        ["Tanggal", DATE_STR],
        ["Penyedia", VENDOR_NAME],
    ]
    t = Table(meta, colWidths=[5 * cm, 10 * cm])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (0, -1), LIGHT_HEX),
        ('BOX', (0, 0), (-1, -1), 0.5, GRAY_HEX),
        ('INNERGRID', (0, 0), (-1, -1), 0.25, HexColor('#BFC7D0')),
        ('FONTNAME', (0, 0), (0, -1), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, -1), 11),
        ('LEFTPADDING', (0, 0), (-1, -1), 8),
        ('TOPPADDING', (0, 0), (-1, -1), 6),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
    ]))
    story.append(t)
    story.append(Spacer(1, 1.2 * cm))
    fee_banner_style = ParagraphStyle('fee_banner', fontSize=14, alignment=TA_CENTER,
                                       textColor=GOLD_HEX, fontName='Helvetica-Bold', leading=18)
    story.append(Paragraph(f"Model Pembiayaan: {FEE_LABEL}", fee_banner_style))
    story.append(Spacer(1, 1 * cm))
    story.append(Paragraph('"Foundation geospasial production-grade: dari prototype 1-kabupaten', quote_style))
    story.append(Paragraph('jadi platform scalable yang siap deploy ke 27 kabupaten Jawa Barat."', quote_style))
    story.append(PageBreak())

    # EXECUTIVE SUMMARY
    story.append(Paragraph("RINGKASAN EKSEKUTIF", h1))
    story.append(Paragraph(
        "Sistem Monitoring Satelit yang sudah terdeploy di DPUTR Kab. Bandung berhasil memetakan 1,09 juta "
        "bangunan. Untuk naik level dari MVP ke platform production multi-kabupaten yang mampu menangani "
        "4 juta+ bangunan dengan response sub-detik, dibutuhkan upgrade fondasi geospasial enterprise-grade. "
        "SIPERA adalah paket upgrade tersebut.", body))
    story.append(Spacer(1, 6 * mm))
    story.append(Paragraph("Angka Kunci", h2))
    kv_data = [
        ["Total bangunan diproses (setelah upgrade)", "4,07 juta (3,7× dari MVP)"],
        ["Response time (cached)", "<300 ms (dari 22-30 detik)"],
        ["Kabupaten siap deploy", "Template ready (~1 hari per kab)"],
        ["Potensi PAD direct (10% realistis)", "Rp 383,8 miliar"],
    ]
    t = Table(kv_data, colWidths=[8.5 * cm, 7 * cm])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (0, -1), LIGHT_HEX),
        ('BOX', (0, 0), (-1, -1), 0.75, NAVY_HEX),
        ('INNERGRID', (0, 0), (-1, -1), 0.25, GRAY_HEX),
        ('FONTNAME', (0, 0), (0, -1), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, -1), 11),
        ('LEFTPADDING', (0, 0), (-1, -1), 8),
        ('TOPPADDING', (0, 0), (-1, -1), 8),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 8),
    ]))
    story.append(t)
    story.append(Spacer(1, 6 * mm))
    story.append(Paragraph(
        '<i>Dataset di-upgrade dari Microsoft Footprints (1,09M) ke kombinasi Microsoft + Google Open Buildings '
        '(4,07M), ditopang stack spatial database industri (PostGIS + Martin Vector Tile + Redis).</i>', body))
    story.append(Paragraph(
        '<b><font color="#28A745">Artinya: dengan investasi Rp 110 juta, potensi tambahan PAD direct '
        'mencapai Rp 281 miliar — rasio ROI &gt;2.800× pada skenario realistis.</font></b>', body))
    story.append(PageBreak())

    # 1. LATAR BELAKANG
    story.append(Paragraph("1. LATAR BELAKANG", h1))
    story.append(Paragraph(
        "Sistem Monitoring Satelit Bangunan (proposal pendahulu, PROP-SATELIT-202604-001) telah berhasil "
        "membangun MVP enforcement bangunan tanpa izin di Kab. Bandung dengan 1,09 juta data Microsoft "
        "Footprints. Setelah implementasi, teridentifikasi tiga kebutuhan upgrade fondamental:", body))
    for b in [
        "Cakupan data Microsoft (1,09M) belum cover seluruh bangunan; Google Open Buildings menambah ~3M polygon yang sebelumnya tidak terdeteksi.",
        "Stack MariaDB + Leaflet geojson tidak scalable untuk 4M+ polygon — render bisa makan 22-30 detik (cold).",
        "Sistem belum bisa di-replikasi ke kabupaten lain tanpa rebuild dari nol — tidak ada template multi-tenant.",
    ]:
        story.append(Paragraph(f"• {b}", bullet))
    story.append(Spacer(1, 6 * mm))
    story.append(Paragraph(
        "SIPERA adalah paket upgrade yang membawa fondasi sistem ke level production enterprise: "
        "PostGIS spatial database, Martin Vector Tile Server (industry-standard MVT), Redis caching layer, "
        "server-side aggregation, dan filter-gated lazy rendering. Hasilnya: platform geospasial scalable "
        "yang siap untuk 27 kabupaten Jawa Barat.", body))
    story.append(PageBreak())

    # 2. OUTPUT & OUTCOME
    story.append(Paragraph("2. OUTPUT &amp; OUTCOME", h1))
    story.append(Paragraph("Output — Yang DPUTR Terima", h2))
    for b in [
        "Spatial database PostGIS 16 (Docker, tuned) berisi 4,07 juta polygon bangunan ter-indeks GIST.",
        "Martin Vector Tile Server (MVT) untuk rendering peta level Mapbox/ESRI.",
        "Redis caching layer dengan 24h TTL + auto-prewarm cron harian.",
        "Server-side spatial aggregation endpoint (cluster z6-13) — render &lt;300ms untuk 4M+ polygon.",
        "Filter-gated rendering UX (choropleth default → drill-down per kecamatan/desa).",
        "Lazy rendering: pan debounce, centroid-first, viewport clipping.",
        "13 ranked performance optimization (load test report + before/after benchmark).",
        "Template deploy multi-kabupaten (per-kab config, ~1 hari setup).",
        "Source code + dokumentasi spatial stack + runbook deploy/migrasi.",
    ]:
        story.append(Paragraph(f"• {b}", bullet))
    story.append(Spacer(1, 4 * mm))

    story.append(Paragraph("Outcome — Manfaat Terukur", h2))
    for b in [
        "Coverage data naik 3,7× (1,09M → 4,07M bangunan) — tidak ada blind spot.",
        "User experience sub-detik untuk peta 4M polygon (sebelumnya 22-30 detik cold).",
        "Foundation siap di-extend untuk fitur lanjutan: housing program, slum mapping, RDTR compliance.",
        "Per-kabupaten deployment cost turun drastis — 1 hari setup vs 11 minggu rebuild dari nol.",
        "Open architecture, no vendor lock-in (PostGIS + MVT = standar industri).",
    ]:
        story.append(Paragraph(f"• {b}", bullet))
    story.append(PageBreak())

    # 3. POTENSI PENDAPATAN (Kabupaten Bandung) — data-driven
    def styled_table(data, col_widths, header_align_right_col=None, first_row_header=True):
        """Build a navy-header Table with light row alternation."""
        style = [
            ('FONTSIZE', (0, 0), (-1, -1), 9.5),
            ('BOX', (0, 0), (-1, -1), 0.75, NAVY_HEX),
            ('INNERGRID', (0, 0), (-1, -1), 0.25, GRAY_HEX),
            ('LEFTPADDING', (0, 0), (-1, -1), 5),
            ('RIGHTPADDING', (0, 0), (-1, -1), 5),
            ('TOPPADDING', (0, 0), (-1, -1), 5),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 5),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
        ]
        if first_row_header:
            style += [
                ('BACKGROUND', (0, 0), (-1, 0), NAVY_HEX),
                ('TEXTCOLOR', (0, 0), (-1, 0), white),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [white, HexColor('#F5F7FA')]),
            ]
        if header_align_right_col is not None:
            style.append(('ALIGN', (header_align_right_col, 1), (header_align_right_col, -1), 'RIGHT'))
        t = Table(data, colWidths=col_widths)
        t.setStyle(TableStyle(style))
        return t

    story.append(Paragraph("3. POTENSI PENDAPATAN (Kabupaten Bandung)", h1))
    story.append(Paragraph(
        "<i>Section ini menghitung potensi PAD untuk DPUTR Kab. Bandung setelah implementasi "
        "SIPERA, dengan semua asumsi dijelaskan sumbernya supaya bisa diaudit ulang oleh tim "
        "DPUTR atau auditor independen.</i>", body))
    story.append(Spacer(1, 4 * mm))

    # 3.1 Sumber Data
    story.append(Paragraph("3.1 Sumber Data & Referensi", h2))
    src_data = [["Data Point", "Sumber"]]
    for r in [
        ("Penduduk & wilayah Kab. Bandung", "BPS Kab. Bandung 2023-2024"),
        ("PBG terbit Kab. Bandung", "SIMBG Kementerian PUPR (5.241 record real)"),
        ("Bangunan terdeteksi 1,09M", "Microsoft GlobalMLBuildingFootprints v3 (2024)"),
        ("Bangunan terdeteksi 4,07M", "Google Open Buildings v3 (rilis Januari 2024)"),
        ("Aturan PBG & retribusi", "PP 16/2021 + Perda Kab. Bandung Retribusi PBG"),
        ("Benchmark enforcement digital", "Studi tax compliance Direktorat Jenderal Pajak"),
    ]:
        src_data.append([wp(r[0]), wp(r[1])])
    story.append(styled_table(src_data, [7 * cm, 8.5 * cm]))
    story.append(Spacer(1, 5 * mm))

    # 3.2 Asumsi
    story.append(Paragraph("3.2 Asumsi Kunci & Justifikasi Angka", h2))
    asu_data = [["Asumsi", "Nilai", "Sumber / Dasar"]]
    for r in [
        ("Penduduk Kab. Bandung", "3,82 juta jiwa", "BPS Kab. Bandung 2023"),
        ("Luas wilayah", "1.762 km²", "BPS Kab. Bandung"),
        ("Total kecamatan", "31 kecamatan", "BPS — 270 desa + 10 kelurahan"),
        ("Bangunan terdeteksi MVP", "1,09 juta", "Microsoft Footprints v3 (2024)"),
        ("Bangunan terdeteksi SIPERA", "4,07 juta", "Microsoft + Google Open Buildings v3"),
        ("Filter ≥200 m² tanpa izin (MVP)", "96.322 (8,8%)", "Hasil filter audit data project"),
        ("Estimasi target SIPERA", "~360.000 unit", "Proporsi 8,8% × 4,07M"),
        ("Retribusi rata-rata PBG", "Rp 10,66 juta", "Real data SIMBG 5.241 PBG = Rp 55,86 M"),
        ("Enforcement realistis 10%", "36.000 PBG", "Benchmark tax digital compliance"),
    ]:
        asu_data.append([wp(r[0]), wp(r[1]), wp(r[2])])
    story.append(styled_table(asu_data, [5.5 * cm, 4 * cm, 6 * cm]))
    story.append(Spacer(1, 4 * mm))
    story.append(Paragraph(
        '<i><font color="#5A5A5A">Validasi cross-check: 3,82 juta penduduk ÷ 4 jiwa per rumah '
        '= 955.000 unit rumah tangga (range plausible terhadap 4,07M total termasuk '
        'non-residensial).</font></i>', body))
    story.append(Spacer(1, 5 * mm))

    # 3.3 Skenario PAD
    story.append(Paragraph("3.3 Skenario Potensi PAD dari 360.000 Bangunan Target", h2))
    scen_data = [["Skenario", "Tingkat", "PBG Tertarik", "Potensi PAD"]]
    for r in [
        ("Super konservatif", "1%", "3.600", "Rp 38,4 miliar"),
        ("Konservatif", "5%", "18.000", "Rp 191,9 miliar"),
        ("Realistis", "10%", "36.000", "Rp 383,8 miliar"),
        ("Target wajar", "20%", "72.000", "Rp 767,5 miliar"),
        ("Agresif", "30%", "108.000", "Rp 1.151 miliar"),
    ]:
        scen_data.append([wp(r[0]), r[1], r[2], r[3]])
    t = styled_table(scen_data, [4.5 * cm, 2.5 * cm, 3.5 * cm, 5 * cm])
    story.append(t)
    story.append(Spacer(1, 4 * mm))
    story.append(Paragraph(
        '<i><font color="#5A5A5A">Formula: PBG tertarik = 360.000 × tingkat enforcement. '
        'Potensi PAD = PBG tertarik × Rp 10,66 juta. Skenario realistis (10%) sebanding dengan '
        'baseline tax compliance improvement DJP setelah digitalisasi (8-15%).</font></i>', body))
    story.append(Spacer(1, 3 * mm))
    story.append(Paragraph(
        '<b><font color="#28A745">Tambahan dari MVP Satelit (10% realistis): '
        'Rp 383,8 M vs Rp 102,7 M = +Rp 281 miliar PAD tambahan yang baru terjangkau '
        'setelah upgrade SIPERA.</font></b>', body))
    story.append(Spacer(1, 5 * mm))

    # 3.4 Justifikasi Tingkat Enforcement
    story.append(Paragraph("3.4 Justifikasi Tingkat Enforcement", h2))
    enf_data = [["Tingkat", "Kondisi", "Basis Argumentasi"]]
    for r in [
        ("1%", "Worst-case adopsi minim", "Tanpa kampanye, no follow-up"),
        ("5%", "Adopsi natural", "Compliance suka rela properti rata-rata"),
        ("10%", "Workflow enforcement jalan", "DJP digital tax improvement 8-15%"),
        ("20%", "Dedicated team enforcement", "Capaian Kota Bekasi/Tangerang"),
        ("30%", "Kampanye + sanksi penuh", "Best practice pemkot dengan SK Bupati"),
    ]:
        enf_data.append([r[0], wp(r[1]), wp(r[2])])
    story.append(styled_table(enf_data, [2 * cm, 5.5 * cm, 8 * cm]))
    story.append(Spacer(1, 5 * mm))

    # 3.5 Productivity Gain
    story.append(Paragraph("3.5 Productivity Gain Staff Enforcement", h2))
    prod_data = [["Aktivitas", "Sebelum", "Setelah SIPERA"]]
    for r in [
        ("Audit 1 kecamatan", "15-20 menit", "5-8 menit (3× cepat)"),
        ("Buka dashboard", "22-30 detik", "<300 ms"),
        ("Pan/zoom peta", "Janky 15 fps", "Smooth 50 fps"),
        ("User concurrent", "5-10 staff", "50-80 staff"),
        ("Akses HP lapangan", "Crash", "Lancar"),
    ]:
        prod_data.append([wp(r[0]), wp(r[1]), wp(r[2])])
    story.append(styled_table(prod_data, [5 * cm, 5 * cm, 5.5 * cm]))
    story.append(Spacer(1, 4 * mm))
    story.append(Paragraph(
        '<i><font color="#5A5A5A">Asumsi: 31 kecamatan × 8 audit/bulan × 12 bulan × 15 menit '
        'hemat = ~750 jam staff/tahun untuk seluruh DPUTR.</font></i>', body))
    story.append(Spacer(1, 5 * mm))

    # 3.6 ROI
    story.append(Paragraph("3.6 Ringkasan ROI SIPERA (Kabupaten Bandung)", h2))
    roi_rows = [
        ("Investasi SIPERA (Paket A+B)", "Rp 100 juta"),
        ("Tambahan PAD direct (skenario realistis 10%)", "+Rp 281 miliar"),
        ("Rasio ROI direct", ">2.800×"),
        ("Break-even", "10 PBG baru ditertibkan"),
        ("Total ROI gabungan (Satelit+SIPERA Rp 260jt)", "~1.476× terhadap PAD Rp 383,8 M"),
    ]
    roi_data = [[wp(k, bold=True), wp(v)] for k, v in roi_rows]
    t = Table(roi_data, colWidths=[9 * cm, 6.5 * cm])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (0, -1), LIGHT_HEX),
        ('BOX', (0, 0), (-1, -1), 0.75, NAVY_HEX),
        ('INNERGRID', (0, 0), (-1, -1), 0.25, GRAY_HEX),
        ('FONTSIZE', (0, 0), (-1, -1), 10),
        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
        ('LEFTPADDING', (0, 0), (-1, -1), 7),
        ('TOPPADDING', (0, 0), (-1, -1), 6),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
    ]))
    story.append(t)
    story.append(Spacer(1, 4 * mm))
    story.append(Paragraph(
        '<i><font color="#1C2D54">Semua angka di atas adalah tambahan dari proposal Satelit '
        'Rp 150 juta. Total ROI dihitung terhadap PAD skenario realistis Kab. Bandung saja '
        '(belum termasuk denda administratif, PBB, atau efek tata ruang).</font></i>', body))
    story.append(PageBreak())

    # 4. RUANG LINGKUP
    story.append(Paragraph("4. RUANG LINGKUP PEKERJAAN", h1))
    story.append(Paragraph(
        "Seluruh scope pekerjaan disajikan dalam bentuk tabel untuk kemudahan audit dan tracking deliverable. "
        "Rincian nilai investasi ada pada Bagian 5.", body))
    story.append(Spacer(1, 4 * mm))

    def pkg_as_table(code, name, items):
        data = [
            ["Kode Paket", code],
            ["Nama Paket", name],
            ["Deliverable & Fitur", ""],
        ]
        for i, it in enumerate(items):
            data.append([f"{code}.{i+1}", wp(it)])
        t = Table(data, colWidths=[3.5 * cm, 12 * cm])
        style = [
            ('BACKGROUND', (0, 0), (0, 1), NAVY_HEX),
            ('TEXTCOLOR', (0, 0), (0, 1), white),
            ('FONTNAME', (0, 0), (0, 1), 'Helvetica-Bold'),
            ('BACKGROUND', (0, 2), (-1, 2), GOLD_HEX),
            ('TEXTCOLOR', (0, 2), (-1, 2), white),
            ('FONTNAME', (0, 2), (-1, 2), 'Helvetica-Bold'),
            ('SPAN', (0, 2), (-1, 2)),
            ('BACKGROUND', (0, 3), (0, -1), LIGHT_HEX),
            ('FONTNAME', (0, 3), (0, -1), 'Helvetica-Bold'),
            ('TEXTCOLOR', (0, 3), (0, -1), NAVY_HEX),
            ('BOX', (0, 0), (-1, -1), 0.75, NAVY_HEX),
            ('INNERGRID', (0, 0), (-1, -1), 0.25, GRAY_HEX),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ('ALIGN', (0, 0), (0, -1), 'CENTER'),
            ('LEFTPADDING', (0, 0), (-1, -1), 6),
            ('RIGHTPADDING', (0, 0), (-1, -1), 6),
            ('TOPPADDING', (0, 0), (-1, -1), 5),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 5),
            ('FONTSIZE', (0, 0), (-1, -1), 10.5),
        ]
        t.setStyle(TableStyle(style))
        story.append(t)
        story.append(Spacer(1, 6 * mm))

    # Paket A
    story.append(Paragraph("Paket A — Foundation Upgrade: Migration & Spatial Stack Setup", h2))
    pkg_as_table("A", "Migrasi MariaDB → PostGIS + Setup Spatial Infrastructure", [
        "Setup PostGIS 16 (Docker) + tuning shared_buffers=2GB, work_mem=64MB, maintenance=512MB.",
        "GIST spatial indexing untuk 4M+ geometries (rendering p95 <500ms).",
        "Migrasi data MariaDB → PostGIS dengan MultiPolygon + Polygon dual-handling.",
        "Setup Redis caching layer (24h TTL + auto-invalidation on verify).",
        "Daily prewarm cron (Windows Task Scheduler / systemd).",
        "Reproducible deploy script (Docker Compose) + rollback runbook.",
    ])

    # Paket B big table
    story.append(Paragraph("Paket B — Geospatial Engineering Platform", h2))
    story.append(Paragraph("Dipecah menjadi 7 sub-paket teknis:", body))
    story.append(Spacer(1, 3 * mm))

    subpkg = [
        ("B.1", "Multi-Source Data Pipeline (4M+)", [
            "Ingest Google Open Buildings 4,07 juta polygon (3,7× dari Microsoft).",
            "Multi-source merger: Microsoft + Google + manual annotation.",
            "Kecamatan enrichment 4,07M rows via Python + scipy KDTree.",
            "High-performance spatial matching engine (~36K rows/sec).",
            "Data quality audit + reconciliation report multi-source.",
            "Reproducible pipeline scripts (Python + Bash + Artisan).",
        ]),
        ("B.2", "Enterprise Spatial DB (PostGIS+Redis)", [
            "PostGIS 16 schema dengan multi-source tables + indexed geometry.",
            "Custom MVT SQL function (zoom-aware, modulo sampling).",
            "Redis cache layer 24h TTL untuk tile + cluster endpoints.",
            "MultiPolygon ST_GeometryN handling untuk irregular shapes.",
            "PostgreSQL tuning + benchmark report (cold vs warm).",
        ]),
        ("B.3", "Vector Tile Pipeline (Mapbox-grade)", [
            "Martin Vector Tile Server (industry-standard MVT).",
            "TilesController endpoint dengan tile cache headers (max-age=86400).",
            "Leaflet.VectorGrid integration (replace GeoJSON untuk 4M polygon).",
            "5-state polygon palette aligned dengan cluster pins.",
            "Modulo sampling (b.id % 5) untuk low-zoom supaya render cepat.",
        ]),
        ("B.4", "Server-Side Spatial Aggregation", [
            "Endpoint /clusters dengan PostGIS grid bucketing (z6-13).",
            "Adaptive zoom-based aggregation (step = 112.5 / 2^zoom).",
            "24h cluster cache + daily prewarm via cron.",
            "Frontend mode-switch otomatis (cluster ↔ vector tile ↔ individual).",
        ]),
        ("B.5", "Filter-Gated Rendering UX", [
            "Choropleth default: 31 polygon kecamatan dengan warna density.",
            "Materialized view kecamatan_violation_stats (daily refresh).",
            "Mode otomatis: no filter → choropleth, kecamatan → cluster+tile, desa → individual.",
            "Empty-state UX yang informatif (bukan layar kosong).",
        ]),
        ("B.6", "Lazy Rendering Optimizations", [
            "Pan/zoom debounce 300ms.",
            "Centroid-first rendering z14-15 (titik dulu, polygon z≥16).",
            "Progressive viewport clipping.",
            "Loading skeleton + abort-on-pan race handling.",
        ]),
        ("B.7", "Performance Engineering & Production", [
            "13 ranked performance optimizations + load test + before/after report.",
            "Playwright-based benchmark dengan network timing breakdown.",
            "Tile feature LIMIT 10K + zoom-conditional sampling.",
            "Verified visual regression: warna polygon/cluster alignment.",
            "Health-check endpoint + monitoring runbook.",
        ]),
    ]

    b_data = [["Kode", "Sub-Paket / Deliverable"]]
    gold_rows = []
    for code, name, items in subpkg:
        gold_rows.append(len(b_data))
        b_data.append([code, name])
        for j, it in enumerate(items):
            b_data.append([f"{code}.{j+1}", wp(it)])

    t = Table(b_data, colWidths=[2.5 * cm, 13 * cm])
    style = [
        ('BACKGROUND', (0, 0), (-1, 0), NAVY_HEX),
        ('TEXTCOLOR', (0, 0), (-1, 0), white),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('BOX', (0, 0), (-1, -1), 0.75, NAVY_HEX),
        ('INNERGRID', (0, 0), (-1, -1), 0.25, GRAY_HEX),
        ('FONTSIZE', (0, 0), (-1, -1), 10),
        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
        ('ALIGN', (0, 0), (0, -1), 'CENTER'),
        ('LEFTPADDING', (0, 0), (-1, -1), 5),
        ('RIGHTPADDING', (0, 0), (-1, -1), 5),
        ('TOPPADDING', (0, 0), (-1, -1), 5),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 5),
    ]
    for gr in gold_rows:
        style.append(('BACKGROUND', (0, gr), (-1, gr), GOLD_HEX))
        style.append(('TEXTCOLOR', (0, gr), (-1, gr), white))
        style.append(('FONTNAME', (0, gr), (-1, gr), 'Helvetica-Bold'))
        style.append(('FONTSIZE', (0, gr), (-1, gr), 11))
    deliv_rows = [i for i in range(1, len(b_data)) if i not in gold_rows]
    for dr in deliv_rows:
        style.append(('BACKGROUND', (0, dr), (0, dr), LIGHT_HEX))
        style.append(('FONTNAME', (0, dr), (0, dr), 'Helvetica-Bold'))
        style.append(('TEXTCOLOR', (0, dr), (0, dr), NAVY_HEX))
    t.setStyle(TableStyle(style))
    story.append(t)
    story.append(Spacer(1, 6 * mm))

    # Paket C & D
    story.append(Paragraph("Paket C — Multi-Kabupaten Templating (Opsional)", h2))
    pkg_as_table("C", "Template Deploy ke Kabupaten Lain (1 hari setup)", [
        "Schema generic (tidak hardcode Kab. Bandung).",
        "Per-kabupaten config file (admin boundary, bbox, prewarm targets).",
        "Documentation deploy ke kabupaten baru (~1 hari setup).",
        "Tested deploy script untuk minimal 1 kabupaten referensi.",
    ])

    story.append(Paragraph("Paket D — Maintenance & SLA 6 Bulan (Opsional)", h2))
    pkg_as_table("D", "Maintenance & Support 6 Bulan", [
        "Bug-fix unlimited untuk defect spatial stack.",
        "Performance regression monitoring bulanan.",
        "Cache + prewarm health check.",
        "Response time 24 jam untuk issue blocker.",
    ])
    story.append(PageBreak())

    # 5. INVESTASI
    story.append(Paragraph("5. INVESTASI &amp; PEMBAYARAN", h1))
    story.append(Paragraph(
        "<i>Setelah melihat output, outcome, dan potensi PAD, berikut rincian investasi "
        "yang dibutuhkan:</i>", body))
    story.append(Spacer(1, 4 * mm))

    contrast_data = [
        ["Potensi PAD (realistis 10%)", "Rp 383,8 miliar"],
        [f"Biaya SIPERA ({'Annual' if FEE_MODE == 'annual' else 'One-Time'})", rupiah(PRICE_TOTAL, with_suffix=True)],
    ]
    t = Table(contrast_data, colWidths=[7 * cm, 8 * cm])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), HexColor('#28A745')),
        ('BACKGROUND', (0, 1), (-1, 1), NAVY_HEX),
        ('TEXTCOLOR', (0, 0), (-1, -1), white),
        ('FONTNAME', (0, 0), (-1, -1), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, -1), 13),
        ('BOX', (0, 0), (-1, -1), 0.75, NAVY_HEX),
        ('INNERGRID', (0, 0), (-1, -1), 0.25, white),
        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
        ('LEFTPADDING', (0, 0), (-1, -1), 10),
        ('TOPPADDING', (0, 0), (-1, -1), 10),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 10),
    ]))
    story.append(t)
    story.append(Spacer(1, 6 * mm))
    story.append(Paragraph(
        '<b><font color="#28A745">ROI direct &gt;2.800× · Break-even: cukup 10 PBG baru ditertibkan '
        'dari 360.000 bangunan terdeteksi.</font></b>', center))
    story.append(Spacer(1, 6 * mm))

    story.append(Paragraph("Rincian Harga per Paket", h2))
    price_data = [
        ["Paket", "Deskripsi", "Harga"],
        ["A", "Foundation Upgrade: Migration & Spatial Stack", rupiah(PRICE_A)],
        ["B", "Geospatial Engineering Platform (7 sub-paket)", rupiah(PRICE_B)],
        ["C", "Multi-Kabupaten Templating (opsional)", rupiah(PRICE_C)],
        ["D", "Maintenance & SLA 6 Bulan (opsional)", rupiah(PRICE_D)],
        ["", f"TOTAL (A+B+C+D){' / tahun' if FEE_MODE == 'annual' else ''}", rupiah(PRICE_TOTAL, with_suffix=True)],
    ]
    t = Table(price_data, colWidths=[2 * cm, 9 * cm, 4 * cm])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), NAVY_HEX),
        ('TEXTCOLOR', (0, 0), (-1, 0), white),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('BACKGROUND', (0, -1), (-1, -1), GOLD_HEX),
        ('FONTNAME', (0, -1), (-1, -1), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, -1), 11),
        ('BOX', (0, 0), (-1, -1), 0.75, NAVY_HEX),
        ('INNERGRID', (0, 0), (-1, -1), 0.25, GRAY_HEX),
        ('ALIGN', (0, 0), (0, -1), 'CENTER'),
        ('ALIGN', (2, 1), (2, -1), 'RIGHT'),
        ('LEFTPADDING', (0, 0), (-1, -1), 6),
        ('TOPPADDING', (0, 0), (-1, -1), 7),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 7),
    ]))
    story.append(t)
    story.append(Spacer(1, 8 * mm))

    story.append(Paragraph("Term Pembayaran", h2))
    if FEE_MODE == 'annual':
        pay_data = [
            ["Milestone", "%", "Nominal"],
            ["Pembayaran Tahun 1 (kickoff + onboarding)", "100%", rupiah(PRICE_TOTAL)],
            ["Renewal Tahun 2", "100%", rupiah(PRICE_TOTAL)],
            ["Renewal Tahun 3 (opsional)", "100%", rupiah(PRICE_TOTAL)],
        ]
    else:
        pay_data = [
            ["Milestone", "%", "Nominal"],
            ["Down Payment (kickoff)", "30%", rupiah(int(PRICE_TOTAL * 0.3))],
            ["Paket A selesai (PostGIS+Redis live)", "20%", rupiah(int(PRICE_TOTAL * 0.2))],
            ["Paket B.1-B.4 selesai (MVT + cluster live)", "30%", rupiah(int(PRICE_TOTAL * 0.3))],
            ["Handover + UAT + benchmark report", "20%", rupiah(int(PRICE_TOTAL * 0.2))],
        ]
    t = Table(pay_data, colWidths=[9 * cm, 2 * cm, 4 * cm])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), NAVY_HEX),
        ('TEXTCOLOR', (0, 0), (-1, 0), white),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, -1), 11),
        ('BOX', (0, 0), (-1, -1), 0.75, NAVY_HEX),
        ('INNERGRID', (0, 0), (-1, -1), 0.25, GRAY_HEX),
        ('ROWBACKGROUNDS', (0, 1), (-1, -1), [white, HexColor('#F5F7FA')]),
        ('ALIGN', (1, 0), (-1, -1), 'CENTER'),
        ('ALIGN', (2, 1), (2, -1), 'RIGHT'),
    ]))
    story.append(t)
    story.append(Spacer(1, 6 * mm))

    notes_block = [Paragraph("Catatan:", h3)]
    base_notes_pdf = [
        "Harga sudah termasuk source code + dokumentasi spatial stack + runbook deploy.",
        "Stand-alone proposal — bisa di-bundle dengan proposal Satelit (Rp 150jt) dengan diskon sinergi 5-10jt.",
        "Deployment di infrastruktur klien (tidak include biaya server/domain).",
        "Penawaran berlaku 30 hari sejak tanggal diterbitkan.",
    ]
    if FEE_MODE == 'annual':
        notes_pdf = base_notes_pdf[:3] + [
            "Annual subscription mencakup: maintenance unlimited bug-fix, performance regression monitoring, security update, minor enhancement (jam kerja terbatas).",
            "Auto-renewal opsional tiap tahun; bisa di-cancel dengan notice 60 hari sebelum periode habis.",
            "SLA response 24 jam untuk issue blocker; uptime + cache health check inklusif.",
        ] + base_notes_pdf[3:]
    else:
        notes_pdf = base_notes_pdf[:3] + ["Garansi bug-fix 60 hari setelah handover untuk spatial stack."] + base_notes_pdf[3:]
    for b in notes_pdf:
        notes_block.append(Paragraph(f"• {b}", bullet))
    story.append(KeepTogether(notes_block))
    story.append(PageBreak())

    # 6. TIMELINE
    story.append(Paragraph("6. TIMELINE", h1))
    timeline_rows = [
        ("Foundation Upgrade (PostGIS + Redis)", "1", "Spatial DB + cache layer live."),
        ("Data Pipeline 4M (Google + matching)", "2", "4,07M row terklasifikasi + indexed."),
        ("Vector Tile Server (Martin MVT)", "3", "MVT endpoint live + Leaflet.VectorGrid."),
        ("Server-Side Clustering", "4", "/clusters endpoint + mode-switch frontend."),
        ("Filter-Gated + Lazy Rendering", "5", "Choropleth + lazy UX."),
        ("Performance Engineering 13 ranks", "6", "Load test report + before/after benchmark."),
        ("Multi-Kab + UAT + Handover", "7", "Template deploy + dokumentasi + sign-off."),
    ]
    timeline_data = [["Fase", "Minggu", "Deliverable"]]
    for r in timeline_rows:
        timeline_data.append([wp(r[0]), r[1], wp(r[2])])
    t = Table(timeline_data, colWidths=[6 * cm, 2 * cm, 7 * cm])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), NAVY_HEX),
        ('TEXTCOLOR', (0, 0), (-1, 0), white),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, -1), 10.5),
        ('BOX', (0, 0), (-1, -1), 0.75, NAVY_HEX),
        ('INNERGRID', (0, 0), (-1, -1), 0.25, GRAY_HEX),
        ('ROWBACKGROUNDS', (0, 1), (-1, -1), [white, HexColor('#F5F7FA')]),
        ('ALIGN', (1, 1), (1, -1), 'CENTER'),
        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
        ('LEFTPADDING', (0, 0), (-1, -1), 6),
        ('TOPPADDING', (0, 0), (-1, -1), 7),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 7),
    ]))
    story.append(t)
    story.append(Spacer(1, 6 * mm))
    story.append(Paragraph('<b><font color="#1C2D54">Total durasi: 7 minggu (~1,75 bulan).</font></b>', body))
    story.append(Paragraph(
        '<i><font color="#5A5A5A">Lebih singkat dari proposal satelit (11 minggu) karena infrastruktur '
        'dasar sudah tersedia dari proyek pendahulu.</font></i>', body))
    story.append(PageBreak())

    # 7. PENUTUP
    story.append(Paragraph("7. PENUTUP", h1))
    story.append(Paragraph(
        "Proposal SIPERA menawarkan upgrade fondasi geospasial dari MVP enforcement (proposal satelit) "
        "ke platform production multi-kabupaten yang scalable. Dengan investasi Rp 110 juta, potensi "
        "tambahan PAD direct Kab. Bandung mencapai Rp 281 miliar (skenario realistis 10%) — ROI "
        "minimum 2.800× pada skenario realistis.", body))
    story.append(Spacer(1, 4 * mm))
    story.append(Paragraph(
        "SIPERA bisa berdiri sendiri sebagai produk baru, atau di-bundle dengan proposal Sistem Monitoring "
        "Satelit sebagai paket lengkap. Output paket ini juga reusable untuk produk turunan: Sistem "
        "Monitoring Perumahan Subsidi, Slum Mapping, Tax Mapping, dan inisiatif geospasial DPUTR lainnya.", body))
    story.append(Spacer(1, 8 * mm))
    story.append(Paragraph("<i>Terima kasih atas kesempatan ini.</i>", center))
    story.append(Spacer(1, 2 * cm))

    sig = [
        [f"{PJ_ROLE},", f"{CLIENT_NAME},"],
        ["", ""],
        ["", ""],
        ["", ""],
        [PJ_NAME, "________________________"],
        ["", ""],
        [f"{EXEC_ROLE},", ""],
        ["", ""],
        ["", ""],
        ["", ""],
        [EXEC_NAME, ""],
    ]
    t = Table(sig, colWidths=[8 * cm, 8 * cm])
    t.setStyle(TableStyle([
        ('FONTSIZE', (0, 0), (-1, -1), 11),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
        ('TOPPADDING', (0, 0), (-1, -1), 4),
        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
        ('VALIGN', (0, 0), (-1, -1), 'TOP'),
    ]))
    story.append(t)

    doc.build(story)
    print(f"✓ PDF:  {out}")
    return out


# ====================== PPTX ======================
def build_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

    NAVY_RGB = RGBColor(28, 45, 84); GOLD_RGB = RGBColor(201, 162, 39)
    WHITE_RGB = RGBColor(255, 255, 255); GRAY_RGB = RGBColor(90, 90, 90)
    SUCCESS_RGB = RGBColor(40, 167, 69); LIGHT_RGB = RGBColor(232, 236, 239)

    blank_layout = prs.slide_layouts[6]

    def add_bg(slide, color=WHITE_RGB):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid(); bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        return bg

    def add_bar(slide, x, y, w, h, color):
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        bar.fill.solid(); bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        return bar

    def add_text(slide, text, x, y, w, h, size=18, bold=False, color=NAVY_RGB,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font='Calibri'):
        tx = slide.shapes.add_textbox(x, y, w, h)
        tf = tx.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
        tf.margin_left = Emu(0); tf.margin_right = Emu(0); tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = font
        return tx

    def add_multiline(slide, lines, x, y, w, h, size=14, bold=False, color=NAVY_RGB,
                      align=PP_ALIGN.LEFT, font='Calibri'):
        tx = slide.shapes.add_textbox(x, y, w, h)
        tf = tx.text_frame; tf.word_wrap = True
        tf.margin_left = Emu(0); tf.margin_right = Emu(0); tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            r = p.add_run(); r.text = line
            r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = font
        return tx

    def add_header(slide, title, num=None):
        add_bar(slide, 0, 0, prs.slide_width, Inches(0.8), NAVY_RGB)
        if num:
            add_text(slide, num, Inches(0.5), Inches(0.2), Inches(1), Inches(0.4),
                     size=14, bold=True, color=GOLD_RGB)
        add_text(slide, title, Inches(1.4) if num else Inches(0.5), Inches(0.15),
                 Inches(11), Inches(0.5), size=22, bold=True, color=WHITE_RGB)
        add_bar(slide, 0, Inches(0.8), prs.slide_width, Inches(0.05), GOLD_RGB)

    # ============ SLIDE 1: COVER ============
    s = prs.slides.add_slide(blank_layout)
    add_bg(s, NAVY_RGB)
    add_bar(s, 0, Inches(3.0), prs.slide_width, Inches(0.1), GOLD_RGB)
    add_text(s, "PROPOSAL PENAWARAN", Inches(0.5), Inches(1.8), Inches(12.3), Inches(0.8),
             size=28, bold=True, color=WHITE_RGB, align=PP_ALIGN.CENTER)
    add_text(s, PROJECT_NAME, Inches(0.5), Inches(2.4), Inches(12.3), Inches(0.6),
             size=22, color=GOLD_RGB, align=PP_ALIGN.CENTER, bold=True)
    add_text(s, PROJECT_SUBTITLE, Inches(0.5), Inches(3.15), Inches(12.3), Inches(0.6),
             size=12, color=LIGHT_RGB, align=PP_ALIGN.CENTER)
    add_text(s, f"Untuk {CLIENT_NAME}", Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.5),
             size=18, color=WHITE_RGB, align=PP_ALIGN.CENTER)
    add_text(s, DATE_STR, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.4),
             size=14, color=LIGHT_RGB, align=PP_ALIGN.CENTER)
    add_text(s, f"Disiapkan oleh: {VENDOR_NAME}", Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.3),
             size=11, color=LIGHT_RGB, align=PP_ALIGN.CENTER)

    # ============ SLIDE 2: KENAPA UPGRADE ============
    s = prs.slides.add_slide(blank_layout); add_bg(s)
    add_header(s, "Kenapa Butuh Upgrade SIPERA?", "01")
    add_text(s, "Tiga Gap Fondamental dari Sistem MVP (Satelit)", Inches(0.5), Inches(1.2),
             Inches(12), Inches(0.5), size=20, bold=True)
    gaps = [
        ("3,7×", "Coverage data: Microsoft 1,09M → Google + Microsoft 4,07M polygon", SUCCESS_RGB),
        ("22-30s", "Cold render time stack MariaDB+Leaflet untuk 4M polygon", GRAY_RGB),
        ("Rebuild", "Tidak ada template multi-kabupaten — rebuild dari nol setiap kab.", GOLD_RGB),
    ]
    for i, (num, txt, color) in enumerate(gaps):
        y = Inches(2.2 + i * 1.4)
        add_text(s, num, Inches(0.7), y, Inches(3), Inches(0.9), size=32, bold=True,
                 color=color, align=PP_ALIGN.LEFT)
        add_text(s, txt, Inches(4.0), y + Inches(0.2), Inches(8.5), Inches(0.7),
                 size=15, color=NAVY_RGB)

    # ============ SLIDE 3: SOLUSI STACK ============
    s = prs.slides.add_slide(blank_layout); add_bg(s)
    add_header(s, "Solusi: Spatial Stack Enterprise-Grade", "02")
    add_text(s, "Stack industri-standar (sama dengan Mapbox / ArcGIS) — dengan biaya freelance",
             Inches(0.5), Inches(1.1), Inches(12), Inches(0.5),
             size=14, color=GRAY_RGB, align=PP_ALIGN.CENTER)

    steps = [
        ("🗺️", "DATA", "Google 4,07M\n+ Microsoft 1,09M\nmulti-source merger"),
        ("🛢️", "DB", "PostGIS 16\n+ GIST index\n+ Redis cache"),
        ("⚡", "TILE", "Martin MVT\nvector tile server\nMapbox-grade"),
        ("🎯", "CLUSTER", "Server-side\naggregation z6-13\n<300ms response"),
        ("🎨", "UX", "Filter-gated\n+ lazy rendering\n+ choropleth"),
    ]
    for i, (ico, title, desc) in enumerate(steps):
        x = Inches(0.5 + i * 2.5); y = Inches(2.2)
        card = add_bar(s, x, y, Inches(2.3), Inches(3.6), LIGHT_RGB)
        card.line.color.rgb = NAVY_RGB; card.line.width = Pt(1)
        add_text(s, ico, x, y + Inches(0.3), Inches(2.3), Inches(0.8),
                 size=32, align=PP_ALIGN.CENTER)
        badge = add_bar(s, x + Inches(0.8), y + Inches(1.2), Inches(0.7), Inches(0.4), NAVY_RGB)
        add_text(s, f"0{i+1}", x + Inches(0.8), y + Inches(1.2), Inches(0.7), Inches(0.4),
                 size=14, bold=True, color=WHITE_RGB, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, title, x + Inches(0.1), y + Inches(1.75), Inches(2.1), Inches(0.4),
                 size=14, bold=True, color=NAVY_RGB, align=PP_ALIGN.CENTER)
        add_multiline(s, desc.split("\n"), x + Inches(0.15), y + Inches(2.25),
                       Inches(2.0), Inches(1.3), size=10, color=GRAY_RGB, align=PP_ALIGN.CENTER)

    # ============ SLIDE 4: FITUR UTAMA ============
    s = prs.slides.add_slide(blank_layout); add_bg(s)
    add_header(s, "Fitur Utama SIPERA", "03")
    fitur_cards = [
        ("Multi-Source Data 4M", "• Google Open Buildings 4,07M\n• Microsoft merger\n• KDTree matching 36K/s\n• Quality audit"),
        ("PostGIS + Redis", "• PostGIS 16 Docker\n• GIST spatial index\n• Redis 24h cache\n• Auto-prewarm cron"),
        ("Vector Tile Pipeline", "• Martin MVT server\n• Leaflet.VectorGrid\n• Modulo sampling\n• Tile cache headers"),
        ("Server-Side Cluster", "• PostGIS grid bucketing\n• Adaptive zoom step\n• <300ms response\n• Daily prewarm"),
        ("Filter-Gated UX", "• Choropleth default\n• Drill-down per kec\n• Empty-state UX\n• Materialized view"),
        ("Lazy Rendering", "• 300ms debounce\n• Centroid-first z14-15\n• Viewport clip\n• 13 perf ranks"),
    ]
    for i, (title, desc) in enumerate(fitur_cards):
        col = i % 3; row = i // 3
        x = Inches(0.4 + col * 4.3); y = Inches(1.3 + row * 2.8)
        card = add_bar(s, x, y, Inches(4.1), Inches(2.6), LIGHT_RGB)
        card.line.color.rgb = NAVY_RGB; card.line.width = Pt(1)
        band = add_bar(s, x, y, Inches(4.1), Inches(0.55), NAVY_RGB)
        add_text(s, title, x, y, Inches(4.1), Inches(0.55), size=15, bold=True,
                 color=WHITE_RGB, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_multiline(s, desc.split("\n"), x + Inches(0.2), y + Inches(0.7),
                       Inches(3.8), Inches(1.8), size=12, color=NAVY_RGB)

    # ============ SLIDE 5: KOMPARASI HARGA PASAR ============
    s = prs.slides.add_slide(blank_layout); add_bg(s)
    add_header(s, "Komparasi Harga Pasar", "04")
    add_text(s, "Scope setara — vendor lain vs SIPERA", Inches(0.5), Inches(1.1),
             Inches(12), Inches(0.5), size=14, color=GRAY_RGB, align=PP_ALIGN.CENTER)

    # Comparison bars (visual price comparison)
    comparisons = [
        ("ESRI ArcGIS Enterprise", "Rp 1.500 jt", 1500, GRAY_RGB),
        ("Telkom Sigma SmartCity", "Rp 1.800 jt", 1800, GRAY_RGB),
        ("MAPID Vendor Lokal", "Rp 500 jt", 500, GRAY_RGB),
        ("Boutique Studio Mid", "Rp 380 jt", 380, GRAY_RGB),
        ("SIPERA (penawaran)", "Rp 110 jt", 110, SUCCESS_RGB),
    ]
    max_price = 1800
    bar_max_w = Inches(7.5)
    for i, (name, price_str, price, color) in enumerate(comparisons):
        y = Inches(2.0 + i * 0.85)
        add_text(s, name, Inches(0.5), y + Inches(0.1), Inches(3.5), Inches(0.5),
                 size=13, color=NAVY_RGB, bold=True)
        bar_w = Inches(7.5 * price / max_price)
        add_bar(s, Inches(4.2), y + Inches(0.05), bar_w, Inches(0.55), color)
        add_text(s, price_str, Inches(4.2) + bar_w + Inches(0.2), y + Inches(0.1),
                 Inches(2), Inches(0.5), size=14, bold=True, color=color)

    add_bar(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.8), SUCCESS_RGB)
    add_text(s, "SIPERA: 3,5× - 16× lebih murah dari vendor komparabel — tanpa lisensi tahunan",
             Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.8),
             size=16, bold=True, color=WHITE_RGB, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ============ SLIDE 6: PRICING TIER ============
    s = prs.slides.add_slide(blank_layout); add_bg(s)
    add_header(s, "Pricing per Tier Vendor", "05")
    add_text(s, "Konteks: di tier mana SIPERA berada", Inches(0.5), Inches(1.1),
             Inches(12), Inches(0.5), size=14, color=GRAY_RGB, align=PP_ALIGN.CENTER)

    tiers = [
        ("TIER 1", "Freelance Solo", "1 dev · no SLA · garansi 60 hari", "Rp 95-130 jt", "◄ SIPERA", SUCCESS_RGB),
        ("TIER 2", "Boutique Studio", "5-15 orang · PM+QA · SLA 1 hari", "Rp 240-320 jt", "", GRAY_RGB),
        ("TIER 3", "GIS Vendor Lokal", "Tim 6-10 · ISO · SLA jam", "Rp 520-720 jt", "", GRAY_RGB),
        ("TIER 4", "Enterprise Multinational", "Tim 15-30 · lisensi vendor", "Rp 1,4 M - 2,5 M", "", NAVY_RGB),
    ]
    for i, (tier, name, desc, price, badge, color) in enumerate(tiers):
        y = Inches(1.9 + i * 1.15)
        # Tier badge
        add_bar(s, Inches(0.5), y, Inches(1.5), Inches(0.95), color)
        add_text(s, tier, Inches(0.5), y, Inches(1.5), Inches(0.95),
                 size=14, bold=True, color=WHITE_RGB, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Content card
        add_bar(s, Inches(2.1), y, Inches(10.7), Inches(0.95), LIGHT_RGB)
        add_text(s, name, Inches(2.3), y + Inches(0.08), Inches(4), Inches(0.4),
                 size=15, bold=True, color=NAVY_RGB)
        add_text(s, desc, Inches(2.3), y + Inches(0.45), Inches(6), Inches(0.4),
                 size=11, color=GRAY_RGB)
        add_text(s, price, Inches(8.5), y + Inches(0.1), Inches(2.5), Inches(0.5),
                 size=16, bold=True, color=color, align=PP_ALIGN.RIGHT)
        if badge:
            add_text(s, badge, Inches(11.0), y + Inches(0.3), Inches(1.7), Inches(0.4),
                     size=13, bold=True, color=color)

    # ============ SLIDE 7: INVESTASI ============
    s = prs.slides.add_slide(blank_layout); add_bg(s)
    add_header(s, "Investasi SIPERA", "06")

    pkg_boxes = [
        ("A", "Foundation Upgrade (PostGIS + Redis)", f"Rp {PRICE_A/1_000_000:.0f} jt", NAVY_RGB),
        ("B", "Geospatial Engineering Platform (7 sub-paket)", f"Rp {PRICE_B/1_000_000:.0f} jt", NAVY_RGB),
        ("C", "Multi-Kabupaten Templating (opsional)", f"Rp {PRICE_C/1_000_000:.0f} jt", GRAY_RGB),
        ("D", "Maintenance 6 Bulan (opsional)", f"Rp {PRICE_D/1_000_000:.0f} jt", GRAY_RGB),
    ]
    for i, (code, desc, price, color) in enumerate(pkg_boxes):
        y = Inches(1.3 + i * 0.9)
        add_bar(s, Inches(0.5), y, Inches(0.9), Inches(0.7), color)
        add_text(s, code, Inches(0.5), y, Inches(0.9), Inches(0.7),
                 size=22, bold=True, color=WHITE_RGB, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, desc, Inches(1.6), y + Inches(0.1), Inches(8), Inches(0.5),
                 size=15, color=NAVY_RGB, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, price, Inches(9.8), y + Inches(0.1), Inches(3), Inches(0.5),
                 size=18, bold=True, color=color, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    y = Inches(5.2)
    add_bar(s, Inches(0.5), y, Inches(12.3), Inches(1.2), GOLD_RGB)
    add_text(s, "TOTAL INVESTASI (A+B+C+D)", Inches(0.5), y, Inches(6), Inches(1.2),
             size=20, bold=True, color=WHITE_RGB, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    pptx_total_suffix = " / tahun" if FEE_MODE == 'annual' else ""
    add_text(s, f"Rp {PRICE_TOTAL/1_000_000:.0f} juta{pptx_total_suffix}", Inches(6.5), y, Inches(6), Inches(1.2),
             size=32, bold=True, color=WHITE_RGB, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, "Bandingkan dengan vendor enterprise scope setara: Rp 600 jt - 1,5 M (5-14× lebih murah)",
             Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
             size=13, color=SUCCESS_RGB, align=PP_ALIGN.CENTER, bold=True)

    # ============ SLIDE 8: TIMELINE ============
    s = prs.slides.add_slide(blank_layout); add_bg(s)
    add_header(s, "Timeline 7 Minggu", "07")

    phases = [
        ("Foundation", "M 1", "PostGIS+Redis live", NAVY_RGB),
        ("Data Pipeline 4M", "M 2", "4,07M classified", NAVY_RGB),
        ("Vector Tile", "M 3", "Martin MVT live", NAVY_RGB),
        ("Cluster Aggreg.", "M 4", "/clusters endpoint", NAVY_RGB),
        ("Filter+Lazy UX", "M 5", "Choropleth UX", SUCCESS_RGB),
        ("Perf Engineering", "M 6", "13 ranks done", SUCCESS_RGB),
        ("Multi-Kab+UAT", "M 7", "Template + sign-off", GOLD_RGB),
    ]
    step_w = Inches(1.7)
    for i, (title, week, deliv, color) in enumerate(phases):
        x = Inches(0.5 + i * 1.8); y = Inches(2.5)
        shp = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW if i < len(phases) - 1 else MSO_SHAPE.RECTANGLE,
                                  x, y, step_w, Inches(1.0))
        shp.fill.solid(); shp.fill.fore_color.rgb = color; shp.line.fill.background()
        add_text(s, week, x, y + Inches(0.15), step_w, Inches(0.4),
                 size=13, bold=True, color=WHITE_RGB, align=PP_ALIGN.CENTER)
        add_text(s, title, x, y + Inches(0.55), step_w, Inches(0.4),
                 size=9, color=WHITE_RGB, align=PP_ALIGN.CENTER)
        add_text(s, deliv, x, y + Inches(1.2), step_w, Inches(0.5),
                 size=9, color=GRAY_RGB, align=PP_ALIGN.CENTER, bold=True)

    add_bar(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.2), LIGHT_RGB)
    add_text(s, "7 minggu · 4 payment milestone · Garansi 60 hari · 4× lebih cepat dari rebuild dari nol",
             Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.2),
             size=16, bold=True, color=NAVY_RGB, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ============ SLIDE 9: CTA ============
    s = prs.slides.add_slide(blank_layout); add_bg(s, NAVY_RGB)
    add_bar(s, 0, Inches(0.5), prs.slide_width, Inches(0.1), GOLD_RGB)

    add_text(s, "Siap Upgrade ke Production?", Inches(0.5), Inches(1.2), Inches(12.3), Inches(1),
             size=44, bold=True, color=WHITE_RGB, align=PP_ALIGN.CENTER)
    add_text(s, f"Rp {PRICE_TOTAL/1_000_000:.0f} juta{pptx_total_suffix} untuk fondasi geospasial sebanding ESRI/Mapbox enterprise",
             Inches(0.5), Inches(2.6), Inches(12.3), Inches(0.6),
             size=18, color=GOLD_RGB, align=PP_ALIGN.CENTER)

    stats_x = [Inches(1.0), Inches(5.1), Inches(9.2)]
    stats = [
        ("4,07 jt", "Polygon bangunan diproses"),
        ("5-14×", "Lebih murah dari vendor komparabel"),
        ("7 minggu", "Dari kickoff ke production"),
    ]
    for (x, (big, small)) in zip(stats_x, stats):
        y = Inches(3.5)
        box = add_bar(s, x, y, Inches(3.3), Inches(2.2), WHITE_RGB)
        box.line.color.rgb = GOLD_RGB; box.line.width = Pt(2)
        add_text(s, big, x, y + Inches(0.3), Inches(3.3), Inches(1),
                 size=32, bold=True, color=NAVY_RGB, align=PP_ALIGN.CENTER)
        add_text(s, small, x, y + Inches(1.3), Inches(3.3), Inches(0.8),
                 size=13, color=GRAY_RGB, align=PP_ALIGN.CENTER)

    add_text(s, f"Hubungi {VENDOR_NAME}", Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.5),
             size=16, color=WHITE_RGB, align=PP_ALIGN.CENTER, bold=True)
    add_text(s, "📧 email · 📱 WhatsApp · 🌐 portfolio", Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
             size=12, color=LIGHT_RGB, align=PP_ALIGN.CENTER)

    out = os.path.join(OUT_DIR, OUTPUT_FILENAME_PPTX)
    prs.save(out)
    print(f"✓ PPTX: {out}")
    return out


if __name__ == "__main__":
    docx_path = build_docx()
    pdf_path = build_pdf()
    # PPTX skipped — user builds presentation manually
    print("\nAll SIPERA proposal docs generated in:", OUT_DIR)
